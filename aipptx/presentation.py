"""Presentation class with semantic slide management.

Provides high-level methods for creating presentations with
AI-friendly, intent-based APIs.
"""

from __future__ import annotations

from pathlib import Path
from typing import TYPE_CHECKING, Any

from pptx import Presentation as PptxPresentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import PP_PLACEHOLDER

from .formatting import format_for_py2ppt, parse_content
from .layout import LayoutType

if TYPE_CHECKING:
    from .template import Template


class Presentation:
    """AI-friendly presentation with semantic slide methods.

    This class wraps python-pptx's Presentation to provide high-level,
    intent-based methods for creating slides. Instead of dealing
    with placeholder indices and layout details, you can simply
    call methods like add_title_slide() or add_comparison_slide().

    Example:
        >>> template = Template("corporate.pptx")
        >>> pres = template.create_presentation()
        >>>
        >>> pres.add_title_slide("Q4 Review", "January 2025")
        >>> pres.add_content_slide("Key Points", [
        ...     "Revenue up 20%",
        ...     "New markets opened",
        ...     "Customer satisfaction high"
        ... ])
        >>> pres.save("output.pptx")
    """

    def __init__(self, template: "Template") -> None:
        """Initialize from a Template.

        Use Template.create_presentation() instead of this constructor.
        """
        self._template = template
        self._layouts = template._layouts

        # Create python-pptx presentation from template
        self._pptx = PptxPresentation(template.path)

        # Remove existing slides (keep only layouts/masters)
        while len(self._pptx.slides) > 0:
            rId = self._pptx.slides._sldIdLst[0].rId
            self._pptx.part.drop_rel(rId)
            del self._pptx.slides._sldIdLst[0]

    @property
    def slide_count(self) -> int:
        """Get the current number of slides."""
        return len(self._pptx.slides)

    @property
    def template(self) -> "Template":
        """Get the template this presentation was created from."""
        return self._template

    def _find_layout_by_type(self, layout_type: LayoutType) -> int:
        """Find the best layout index for a layout type."""
        for layout in self._layouts:
            if layout.layout_type == layout_type:
                return layout.index
        # Fall back to first layout
        return 0

    def _find_layout(self, layout: str | int | None, layout_type: LayoutType) -> int:
        """Find layout by name, index, or type."""
        if layout is None or layout == "auto":
            return self._find_layout_by_type(layout_type)

        if isinstance(layout, int):
            return layout

        # Fuzzy name match
        layout_lower = layout.lower()
        for l in self._layouts:
            if layout_lower in l.name.lower() or l.name.lower() in layout_lower:
                return l.index
        return self._find_layout_by_type(layout_type)

    def _get_placeholder(self, slide, ph_type):
        """Get a placeholder by type from a slide."""
        for shape in slide.placeholders:
            if shape.placeholder_format.type == ph_type:
                return shape
        return None

    def _set_text_frame(self, shape, text: str) -> None:
        """Set text in a shape's text frame."""
        if shape is None:
            return
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = text

    def _set_body_content(self, shape, content: list, levels: list[int] | None = None) -> None:
        """Set bullet content in a body placeholder."""
        if shape is None:
            return

        tf = shape.text_frame
        tf.clear()

        for i, item in enumerate(content):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            # Set level
            if levels and i < len(levels):
                p.level = levels[i]

            # Set text
            if isinstance(item, str):
                p.text = item
            elif isinstance(item, list):
                # Rich text - just use plain text for now
                p.text = "".join(
                    seg.get("text", "") if isinstance(seg, dict) else str(seg)
                    for seg in item
                )
            elif isinstance(item, dict):
                p.text = item.get("text", "")

    def add_title_slide(
        self,
        title: str,
        subtitle: str = "",
        *,
        layout: str | int | None = None,
    ) -> int:
        """Add a title/cover slide.

        Args:
            title: Main presentation title
            subtitle: Optional subtitle (presenter name, date, etc.)
            layout: Layout name, index, or None for auto-selection

        Returns:
            Slide number of the new slide

        Example:
            >>> pres.add_title_slide("Q4 Business Review", "January 2025")
        """
        layout_idx = self._find_layout(layout, LayoutType.TITLE)
        slide_layout = self._pptx.slide_layouts[layout_idx]
        slide = self._pptx.slides.add_slide(slide_layout)

        # Set title
        title_ph = self._get_placeholder(slide, PP_PLACEHOLDER.TITLE)
        if title_ph is None:
            title_ph = self._get_placeholder(slide, PP_PLACEHOLDER.CENTER_TITLE)
        self._set_text_frame(title_ph, title)

        # Set subtitle
        if subtitle:
            subtitle_ph = self._get_placeholder(slide, PP_PLACEHOLDER.SUBTITLE)
            self._set_text_frame(subtitle_ph, subtitle)

        return len(self._pptx.slides)

    def add_section_slide(
        self,
        title: str,
        subtitle: str = "",
        *,
        layout: str | int | None = None,
    ) -> int:
        """Add a section header/divider slide.

        Args:
            title: Section title
            subtitle: Optional subtitle or description
            layout: Layout name, index, or None for auto-selection

        Returns:
            Slide number of the new slide

        Example:
            >>> pres.add_section_slide("Part 2: Analysis")
        """
        layout_idx = self._find_layout(layout, LayoutType.SECTION)
        slide_layout = self._pptx.slide_layouts[layout_idx]
        slide = self._pptx.slides.add_slide(slide_layout)

        # Set title
        title_ph = self._get_placeholder(slide, PP_PLACEHOLDER.TITLE)
        self._set_text_frame(title_ph, title)

        # Set subtitle if provided
        if subtitle:
            subtitle_ph = self._get_placeholder(slide, PP_PLACEHOLDER.SUBTITLE)
            if subtitle_ph:
                self._set_text_frame(subtitle_ph, subtitle)

        return len(self._pptx.slides)

    def add_content_slide(
        self,
        title: str,
        content: str | list[str | tuple | dict | list],
        *,
        levels: list[int] | None = None,
        layout: str | int | None = None,
    ) -> int:
        """Add a content slide with bullets.

        Args:
            title: Slide title
            content: Content as string, list of strings, or rich text
            levels: Optional indent levels for each bullet (0=top level)
            layout: Layout name, index, or None for auto-selection

        Returns:
            Slide number of the new slide

        Example:
            >>> # Simple bullets
            >>> pres.add_content_slide("Key Points", [
            ...     "First point",
            ...     "Second point",
            ...     "Third point"
            ... ])

            >>> # With nested levels
            >>> pres.add_content_slide("Details", [
            ...     "Main topic",
            ...     "Sub-point 1",
            ...     "Sub-point 2",
            ...     "Another main topic"
            ... ], levels=[0, 1, 1, 0])
        """
        layout_idx = self._find_layout(layout, LayoutType.CONTENT)
        slide_layout = self._pptx.slide_layouts[layout_idx]
        slide = self._pptx.slides.add_slide(slide_layout)

        # Set title
        title_ph = self._get_placeholder(slide, PP_PLACEHOLDER.TITLE)
        self._set_text_frame(title_ph, title)

        # Parse and format content
        paragraphs = parse_content(content, levels)
        formatted_content, formatted_levels = format_for_py2ppt(paragraphs)

        # Set body
        body_ph = self._get_placeholder(slide, PP_PLACEHOLDER.BODY)
        if body_ph is None:
            body_ph = self._get_placeholder(slide, PP_PLACEHOLDER.OBJECT)
        self._set_body_content(body_ph, formatted_content, formatted_levels)

        return len(self._pptx.slides)

    def add_two_column_slide(
        self,
        title: str,
        left_content: list[str | tuple | dict | list],
        right_content: list[str | tuple | dict | list],
        *,
        left_levels: list[int] | None = None,
        right_levels: list[int] | None = None,
        layout: str | int | None = None,
    ) -> int:
        """Add a two-column content slide.

        Args:
            title: Slide title
            left_content: Content for left column
            right_content: Content for right column
            left_levels: Indent levels for left content
            right_levels: Indent levels for right content
            layout: Layout name, index, or None for auto-selection

        Returns:
            Slide number of the new slide

        Example:
            >>> pres.add_two_column_slide(
            ...     "Topics",
            ...     ["Point A", "Point B"],
            ...     ["Point X", "Point Y"]
            ... )
        """
        layout_idx = self._find_layout(layout, LayoutType.TWO_COLUMN)
        slide_layout = self._pptx.slide_layouts[layout_idx]
        slide = self._pptx.slides.add_slide(slide_layout)

        # Set title
        title_ph = self._get_placeholder(slide, PP_PLACEHOLDER.TITLE)
        self._set_text_frame(title_ph, title)

        # Find body placeholders (there should be two)
        body_phs = [
            shape for shape in slide.placeholders
            if shape.placeholder_format.type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT)
        ]
        body_phs.sort(key=lambda s: s.left)  # Sort by x position

        # Format content
        left_paragraphs = parse_content(left_content, left_levels)
        left_formatted, left_lvls = format_for_py2ppt(left_paragraphs)

        right_paragraphs = parse_content(right_content, right_levels)
        right_formatted, right_lvls = format_for_py2ppt(right_paragraphs)

        # Set left content
        if len(body_phs) >= 1:
            self._set_body_content(body_phs[0], left_formatted, left_lvls)

        # Set right content
        if len(body_phs) >= 2:
            self._set_body_content(body_phs[1], right_formatted, right_lvls)

        return len(self._pptx.slides)

    def add_comparison_slide(
        self,
        title: str,
        left_heading: str,
        left_content: list[str | tuple | dict | list],
        right_heading: str,
        right_content: list[str | tuple | dict | list],
        *,
        layout: str | int | None = None,
    ) -> int:
        """Add a comparison slide with two titled columns.

        Ideal for before/after, pros/cons, or versus comparisons.

        Args:
            title: Slide title
            left_heading: Heading for left column
            left_content: Bullet points for left column
            right_heading: Heading for right column
            right_content: Bullet points for right column
            layout: Layout name, index, or None for auto-selection

        Returns:
            Slide number of the new slide

        Example:
            >>> pres.add_comparison_slide(
            ...     "Before vs After",
            ...     "Legacy System",
            ...     ["Slow", "Manual", "Error-prone"],
            ...     "New Platform",
            ...     ["Fast", "Automated", "Reliable"]
            ... )
        """
        # Try to find a comparison layout, fall back to two-column
        layout_idx = self._find_layout(layout, LayoutType.COMPARISON)
        slide_layout = self._pptx.slide_layouts[layout_idx]
        slide = self._pptx.slides.add_slide(slide_layout)

        # Set title
        title_ph = self._get_placeholder(slide, PP_PLACEHOLDER.TITLE)
        self._set_text_frame(title_ph, title)

        # Find body placeholders
        body_phs = [
            shape for shape in slide.placeholders
            if shape.placeholder_format.type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT)
        ]
        # Sort by position: top row first, then left to right
        body_phs.sort(key=lambda s: (s.top, s.left))

        if len(body_phs) >= 4:
            # True comparison layout: heading, content, heading, content
            self._set_text_frame(body_phs[0], left_heading)
            left_paragraphs = parse_content(left_content)
            left_formatted, left_lvls = format_for_py2ppt(left_paragraphs)
            self._set_body_content(body_phs[2], left_formatted, left_lvls)

            self._set_text_frame(body_phs[1], right_heading)
            right_paragraphs = parse_content(right_content)
            right_formatted, right_lvls = format_for_py2ppt(right_paragraphs)
            self._set_body_content(body_phs[3], right_formatted, right_lvls)

        elif len(body_phs) >= 2:
            # Two-column layout - combine heading with content
            left_combined = [left_heading] + list(left_content)
            right_combined = [right_heading] + list(right_content)

            left_levels = [0] + [1] * len(left_content)
            right_levels = [0] + [1] * len(right_content)

            left_paragraphs = parse_content(left_combined, left_levels)
            right_paragraphs = parse_content(right_combined, right_levels)

            left_formatted, left_lvls = format_for_py2ppt(left_paragraphs)
            right_formatted, right_lvls = format_for_py2ppt(right_paragraphs)

            # Sort by x position for left/right
            body_phs.sort(key=lambda s: s.left)
            self._set_body_content(body_phs[0], left_formatted, left_lvls)
            self._set_body_content(body_phs[1], right_formatted, right_lvls)

        elif len(body_phs) >= 1:
            # Single body - combine all
            combined = [
                left_heading,
                *list(left_content),
                "",
                right_heading,
                *list(right_content),
            ]
            combined_levels = (
                [0] + [1] * len(left_content) + [0] + [0] + [1] * len(right_content)
            )
            paragraphs = parse_content(combined, combined_levels)
            formatted, lvls = format_for_py2ppt(paragraphs)
            self._set_body_content(body_phs[0], formatted, lvls)

        return len(self._pptx.slides)

    def add_image_slide(
        self,
        title: str,
        image_path: str | Path,
        caption: str = "",
        *,
        layout: str | int | None = None,
    ) -> int:
        """Add a slide with an image.

        Args:
            title: Slide title
            image_path: Path to the image file
            caption: Optional caption text
            layout: Layout name, index, or None for auto-selection

        Returns:
            Slide number of the new slide

        Example:
            >>> pres.add_image_slide(
            ...     "Product Photo",
            ...     "product.png",
            ...     "Our flagship product"
            ... )
        """
        layout_idx = self._find_layout(layout, LayoutType.IMAGE_CONTENT)
        slide_layout = self._pptx.slide_layouts[layout_idx]
        slide = self._pptx.slides.add_slide(slide_layout)

        # Set title
        title_ph = self._get_placeholder(slide, PP_PLACEHOLDER.TITLE)
        self._set_text_frame(title_ph, title)

        # Add image
        image_path = Path(image_path)
        if image_path.exists():
            slide.shapes.add_picture(
                str(image_path),
                Inches(1),
                Inches(2),
                width=Inches(5),
            )

        # Set caption if provided
        if caption:
            body_ph = self._get_placeholder(slide, PP_PLACEHOLDER.BODY)
            if body_ph:
                self._set_text_frame(body_ph, caption)

        return len(self._pptx.slides)

    def add_blank_slide(self, layout: str | int | None = None) -> int:
        """Add a blank slide.

        Args:
            layout: Layout name, index, or None for auto-selection

        Returns:
            Slide number of the new slide

        Example:
            >>> slide_num = pres.add_blank_slide()
        """
        layout_idx = self._find_layout(layout, LayoutType.BLANK)
        slide_layout = self._pptx.slide_layouts[layout_idx]
        self._pptx.slides.add_slide(slide_layout)
        return len(self._pptx.slides)

    def add_slide(
        self,
        layout: str | int = "auto",
        content_type: str = "content",
        title: str = "",
        content: list[str] | None = None,
        **kwargs: Any,
    ) -> int:
        """Add a slide with auto-layout selection.

        A flexible method that chooses the right slide type
        based on content_type and delegates to the appropriate
        specialized method.

        Args:
            layout: Layout name, index, or "auto"
            content_type: Type of content ("title", "content", "comparison", etc.)
            title: Slide title
            content: Content items
            **kwargs: Additional arguments for the specific slide type

        Returns:
            Slide number of the new slide

        Example:
            >>> pres.add_slide(
            ...     content_type="bullets",
            ...     title="Key Points",
            ...     content=["Point 1", "Point 2"]
            ... )
        """
        content_lower = content_type.lower()

        if content_lower in ("title", "cover", "opening"):
            subtitle = kwargs.get("subtitle", "")
            return self.add_title_slide(title, subtitle, layout=layout)

        elif content_lower in ("section", "divider"):
            subtitle = kwargs.get("subtitle", "")
            return self.add_section_slide(title, subtitle, layout=layout)

        elif content_lower in ("comparison", "versus", "vs"):
            return self.add_comparison_slide(
                title,
                kwargs.get("left_heading", "Option A"),
                kwargs.get("left_content", content or []),
                kwargs.get("right_heading", "Option B"),
                kwargs.get("right_content", []),
                layout=layout,
            )

        elif content_lower in ("two_column", "split"):
            return self.add_two_column_slide(
                title,
                kwargs.get("left_content", content or []),
                kwargs.get("right_content", []),
                layout=layout,
            )

        elif content_lower in ("image", "picture"):
            return self.add_image_slide(
                title,
                kwargs.get("image_path", ""),
                kwargs.get("caption", ""),
                layout=layout,
            )

        elif content_lower == "blank":
            return self.add_blank_slide(layout=layout)

        else:
            # Default to content slide
            return self.add_content_slide(
                title,
                content or [],
                levels=kwargs.get("levels"),
                layout=layout,
            )

    def set_notes(self, slide_number: int, notes: str) -> None:
        """Set speaker notes for a slide.

        Args:
            slide_number: The slide number (1-indexed)
            notes: Notes text (can include newlines)

        Example:
            >>> pres.set_notes(1, "Remember to mention...")
        """
        if slide_number < 1 or slide_number > len(self._pptx.slides):
            return

        slide = self._pptx.slides[slide_number - 1]
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

    def save(self, path: str | Path) -> None:
        """Save the presentation.

        Args:
            path: Output file path

        Example:
            >>> pres.save("output.pptx")
        """
        self._pptx.save(path)

    def __repr__(self) -> str:
        return f"Presentation({self.slide_count} slides, template={self._template.path.name})"
