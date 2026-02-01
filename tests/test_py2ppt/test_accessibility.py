"""Tests for accessibility and optimization features."""

from pathlib import Path

import pytest
from pptx import Presentation as PptxPresentation

from py2ppt import Template, Presentation
from py2ppt.validation import IssueCategory


@pytest.fixture
def template(tmp_path: Path) -> Template:
    """Create a template for testing."""
    template_path = tmp_path / "template.pptx"
    pres = PptxPresentation()
    pres.slides.add_slide(pres.slide_layouts[0])
    pres.save(str(template_path))
    return Template(template_path)


@pytest.fixture
def presentation(template: Template) -> Presentation:
    """Create a presentation from template."""
    return template.create_presentation()


class TestCheckAccessibility:
    """Tests for check_accessibility method."""

    def test_check_empty_presentation(self, presentation: Presentation) -> None:
        """Test accessibility check on empty presentation."""
        result = presentation.check_accessibility()

        assert result.is_valid
        assert isinstance(result.score, float)

    def test_check_basic_presentation(self, presentation: Presentation) -> None:
        """Test accessibility check on basic presentation."""
        presentation.add_title_slide("Title", "Subtitle")
        presentation.add_content_slide("Content", ["Point 1", "Point 2"])

        result = presentation.check_accessibility()

        assert isinstance(result.score, float)
        assert result.score <= 100

    def test_check_missing_title(self, presentation: Presentation) -> None:
        """Test that missing title is flagged."""
        presentation.add_blank_slide()

        result = presentation.check_accessibility()

        # Should have issues about missing title
        accessibility_issues = [
            i for i in result.issues
            if i.category == IssueCategory.ACCESSIBILITY
        ]
        assert len(accessibility_issues) > 0

    def test_check_with_notes(self, presentation: Presentation) -> None:
        """Test accessibility check with speaker notes."""
        presentation.add_content_slide("Slide", ["Point"])
        presentation.set_notes(1, "Notes for accessibility")

        result = presentation.check_accessibility()

        assert isinstance(result.score, float)


class TestOptimizeSlide:
    """Tests for optimize_slide method."""

    def test_optimize_basic_slide(self, presentation: Presentation) -> None:
        """Test basic slide optimization."""
        presentation.add_content_slide("Title", ["Point 1", "Point 2"])

        changes = presentation.optimize_slide(1)

        assert "slide_number" in changes
        assert changes["slide_number"] == 1
        assert "suggestions" in changes

    def test_optimize_slide_many_bullets(self, presentation: Presentation) -> None:
        """Test optimization suggests splitting for many bullets."""
        bullets = [f"Point {i}" for i in range(10)]
        presentation.add_content_slide("Many Points", bullets)

        changes = presentation.optimize_slide(1)

        assert len(changes["suggestions"]) > 0
        assert any("splitting" in s.lower() for s in changes["suggestions"])

    def test_optimize_slide_long_content(self, presentation: Presentation) -> None:
        """Test optimization catches long content."""
        long_bullet = "This is a very long bullet point that goes on and on " * 5
        presentation.add_content_slide("Long Content", [long_bullet])

        changes = presentation.optimize_slide(1)

        assert len(changes["suggestions"]) > 0

    def test_optimize_invalid_slide(self, presentation: Presentation) -> None:
        """Test optimizing invalid slide raises error."""
        from py2ppt.errors import SlideNotFoundError

        presentation.add_title_slide("Title", "")

        with pytest.raises(SlideNotFoundError):
            presentation.optimize_slide(99)


class TestOptimizeAll:
    """Tests for optimize_all method."""

    def test_optimize_all_basic(self, presentation: Presentation) -> None:
        """Test optimizing all slides."""
        presentation.add_title_slide("Title", "")
        presentation.add_content_slide("Content 1", ["Point"])
        presentation.add_content_slide("Content 2", ["Point"])

        all_changes = presentation.optimize_all()

        assert len(all_changes) == 3
        for change in all_changes:
            assert "slide_number" in change
            assert "suggestions" in change

    def test_optimize_all_empty(self, presentation: Presentation) -> None:
        """Test optimizing empty presentation."""
        all_changes = presentation.optimize_all()

        assert len(all_changes) == 0


class TestImagePlaceholders:
    """Tests for image placeholder features."""

    def test_add_image_placeholder(self, presentation: Presentation) -> None:
        """Test adding an image placeholder."""
        presentation.add_blank_slide()

        ph_id = presentation.add_image_placeholder(
            1, "Professional team meeting",
            left=1, top=2, width=5, height=3
        )

        assert ph_id is not None
        assert "img_placeholder" in ph_id

    def test_get_image_placeholders(self, presentation: Presentation) -> None:
        """Test getting all image placeholders."""
        presentation.add_blank_slide()

        presentation.add_image_placeholder(1, "First image", 1, 1, 2, 2)
        presentation.add_image_placeholder(1, "Second image", 4, 1, 2, 2)

        placeholders = presentation.get_image_placeholders()

        assert len(placeholders) == 2
        assert any("First image" in ph["prompt"] for ph in placeholders)
        assert any("Second image" in ph["prompt"] for ph in placeholders)

    def test_fill_image_placeholder(
        self, presentation: Presentation, tmp_path: Path
    ) -> None:
        """Test filling an image placeholder."""
        presentation.add_blank_slide()

        ph_id = presentation.add_image_placeholder(
            1, "Test image", 1, 1, 3, 2
        )

        # Create a test image
        from PIL import Image

        img_path = tmp_path / "test_image.png"
        img = Image.new("RGB", (100, 100), color="red")
        img.save(img_path)

        presentation.fill_image_placeholder(1, ph_id, img_path)

        # Placeholder should be gone, image should be there
        placeholders = presentation.get_image_placeholders()
        assert len(placeholders) == 0

    def test_fill_placeholder_invalid_image(
        self, presentation: Presentation
    ) -> None:
        """Test filling placeholder with missing image raises error."""
        from py2ppt.errors import InvalidDataError

        presentation.add_blank_slide()
        ph_id = presentation.add_image_placeholder(1, "Test", 1, 1, 2, 2)

        with pytest.raises(InvalidDataError):
            presentation.fill_image_placeholder(1, ph_id, "nonexistent.png")

    def test_fill_placeholder_invalid_id(
        self, presentation: Presentation, tmp_path: Path
    ) -> None:
        """Test filling nonexistent placeholder raises error."""
        from py2ppt.errors import InvalidDataError
        from PIL import Image

        presentation.add_blank_slide()

        img_path = tmp_path / "test.png"
        img = Image.new("RGB", (50, 50), color="blue")
        img.save(img_path)

        with pytest.raises(InvalidDataError):
            presentation.fill_image_placeholder(1, "nonexistent_ph", img_path)


class TestSetAltText:
    """Tests for set_alt_text method."""

    def test_set_alt_text_on_shape(self, presentation: Presentation) -> None:
        """Test setting alt text on a shape."""
        presentation.add_blank_slide()
        shape_name = presentation.add_shape(1, "rectangle", 1, 1, 2, 2)

        presentation.set_alt_text(1, shape_name, "A blue rectangle")

        # Should not raise

    def test_set_alt_text_invalid_shape(self, presentation: Presentation) -> None:
        """Test setting alt text on nonexistent shape raises error."""
        from py2ppt.errors import InvalidDataError

        presentation.add_blank_slide()

        with pytest.raises(InvalidDataError):
            presentation.set_alt_text(1, "NonexistentShape", "Alt text")


class TestValidateWithAccessibility:
    """Tests for validate with accessibility option."""

    def test_validate_include_accessibility(self, presentation: Presentation) -> None:
        """Test validate with include_accessibility flag."""
        presentation.add_blank_slide()

        result = presentation.validate(include_accessibility=True)

        # Should include accessibility issues
        has_accessibility = any(
            i.category == IssueCategory.ACCESSIBILITY
            for i in result.issues
        )
        assert has_accessibility

    def test_validate_with_brand_rules(self, presentation: Presentation) -> None:
        """Test validate with brand rules."""
        presentation.add_content_slide(
            "Title",
            ["Point 1", "Point 2", "Point 3", "Point 4", "Point 5", "Point 6"]
        )

        result = presentation.validate(brand_rules={"max_bullets": 3})

        # Should flag the bullet count
        has_bullet_issue = any(
            "bullets" in i.message.lower()
            for i in result.issues
        )
        assert has_bullet_issue


class TestSaveWithAccessibility:
    """Tests for saving presentations with accessibility features."""

    def test_save_with_placeholders(
        self, presentation: Presentation, tmp_path: Path
    ) -> None:
        """Test saving presentation with image placeholders."""
        presentation.add_blank_slide()
        presentation.add_image_placeholder(1, "Test placeholder", 1, 1, 3, 2)

        output_path = tmp_path / "accessibility_output.pptx"
        presentation.save(output_path)

        assert output_path.exists()

        # Verify file can be opened
        loaded = PptxPresentation(str(output_path))
        assert len(loaded.slides) == 1
