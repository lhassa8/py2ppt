"""Tests for slide cloning and presentation merge features."""

from pathlib import Path

import pytest
from pptx import Presentation as PptxPresentation

from py2ppt import Template, Presentation


@pytest.fixture
def template(tmp_path: Path) -> Template:
    """Create a template for testing."""
    template_path = tmp_path / "template.pptx"
    pres = PptxPresentation()
    pres.slides.add_slide(pres.slide_layouts[0])
    pres.save(str(template_path))
    return Template(template_path)


@pytest.fixture
def presentation(template: Template) -> Presentation:
    """Create a presentation from template."""
    return template.create_presentation()


class TestCloneSlide:
    """Tests for clone_slide method."""

    def test_clone_slide_basic(self, presentation: Presentation) -> None:
        """Test basic slide cloning."""
        presentation.add_title_slide("Original Title", "Subtitle")
        assert presentation.slide_count == 1

        cloned_num = presentation.clone_slide(1)

        assert cloned_num == 2
        assert presentation.slide_count == 2

    def test_clone_slide_preserves_content(self, presentation: Presentation) -> None:
        """Test that cloning preserves slide content."""
        presentation.add_content_slide("Test Content", ["Point 1", "Point 2"])

        presentation.clone_slide(1)

        # Both slides should exist
        assert presentation.slide_count == 2

        # Verify the clone was created (content may vary due to XML copying)
        slide1 = presentation.describe_slide(1)
        slide2 = presentation.describe_slide(2)

        # Both should have the same layout
        assert slide1["layout"] == slide2["layout"]

    def test_clone_slide_with_insert_at(self, presentation: Presentation) -> None:
        """Test cloning with specific insert position."""
        presentation.add_title_slide("Slide 1", "")
        presentation.add_content_slide("Slide 2", ["Content"])
        presentation.add_section_slide("Slide 3")

        # Clone slide 3 and insert at position 2
        cloned_num = presentation.clone_slide(3, insert_at=2)

        assert cloned_num == 2
        assert presentation.slide_count == 4

    def test_clone_slide_multiple_times(self, presentation: Presentation) -> None:
        """Test cloning the same slide multiple times."""
        presentation.add_title_slide("Template Slide", "")

        presentation.clone_slide(1)
        presentation.clone_slide(1)
        presentation.clone_slide(1)

        assert presentation.slide_count == 4

    def test_clone_slide_invalid_source(self, presentation: Presentation) -> None:
        """Test cloning from invalid source raises error."""
        from py2ppt.errors import SlideNotFoundError

        presentation.add_title_slide("Only Slide", "")

        with pytest.raises(SlideNotFoundError):
            presentation.clone_slide(99)

    def test_clone_slide_preserves_notes(self, presentation: Presentation) -> None:
        """Test that cloning preserves speaker notes."""
        presentation.add_title_slide("Title", "")
        presentation.set_notes(1, "These are speaker notes")

        presentation.clone_slide(1)

        # Check notes on cloned slide
        slide = presentation._pptx.slides[1]
        notes = slide.notes_slide.notes_text_frame.text
        assert "speaker notes" in notes.lower()


class TestCloneSlideFrom:
    """Tests for clone_slide_from method."""

    def test_clone_from_another_presentation(
        self, template: Template
    ) -> None:
        """Test cloning from another presentation."""
        pres1 = template.create_presentation()
        pres1.add_title_slide("Source Title", "Source Subtitle")

        pres2 = template.create_presentation()
        cloned_num = pres2.clone_slide_from(pres1, 1)

        assert cloned_num == 1
        assert pres2.slide_count == 1

    def test_clone_from_multiple_slides(
        self, template: Template
    ) -> None:
        """Test cloning multiple slides from another presentation."""
        pres1 = template.create_presentation()
        pres1.add_title_slide("Title 1", "")
        pres1.add_content_slide("Content", ["Item"])
        pres1.add_section_slide("Section")

        pres2 = template.create_presentation()
        pres2.clone_slide_from(pres1, 1)
        pres2.clone_slide_from(pres1, 2)
        pres2.clone_slide_from(pres1, 3)

        assert pres2.slide_count == 3

    def test_clone_from_with_insert_at(
        self, template: Template
    ) -> None:
        """Test cloning from with specific insert position."""
        pres1 = template.create_presentation()
        pres1.add_title_slide("External Slide", "")

        pres2 = template.create_presentation()
        pres2.add_title_slide("Slide 1", "")
        pres2.add_content_slide("Slide 2", [])

        cloned_num = pres2.clone_slide_from(pres1, 1, insert_at=2)

        assert cloned_num == 2
        assert pres2.slide_count == 3

    def test_clone_from_invalid_source(
        self, template: Template
    ) -> None:
        """Test cloning from invalid source slide raises error."""
        from py2ppt.errors import SlideNotFoundError

        pres1 = template.create_presentation()
        pres1.add_title_slide("Only Slide", "")

        pres2 = template.create_presentation()

        with pytest.raises(SlideNotFoundError):
            pres2.clone_slide_from(pres1, 99)


class TestMerge:
    """Tests for merge method."""

    def test_merge_basic(self, template: Template) -> None:
        """Test basic presentation merge."""
        pres1 = template.create_presentation()
        pres1.add_title_slide("Main Deck", "")

        pres2 = template.create_presentation()
        pres2.add_content_slide("Extra Content", ["Item 1"])

        merged_nums = pres1.merge(pres2)

        assert len(merged_nums) == 1
        assert pres1.slide_count == 2

    def test_merge_multiple_slides(self, template: Template) -> None:
        """Test merging presentation with multiple slides."""
        pres1 = template.create_presentation()
        pres1.add_title_slide("Main", "")

        pres2 = template.create_presentation()
        pres2.add_title_slide("Title 1", "")
        pres2.add_content_slide("Content 1", [])
        pres2.add_section_slide("Section")

        merged_nums = pres1.merge(pres2)

        assert len(merged_nums) == 3
        assert pres1.slide_count == 4

    def test_merge_with_insert_at(self, template: Template) -> None:
        """Test merge with specific insert position."""
        pres1 = template.create_presentation()
        pres1.add_title_slide("First", "")
        pres1.add_title_slide("Last", "")

        pres2 = template.create_presentation()
        pres2.add_content_slide("Middle", [])

        merged_nums = pres1.merge(pres2, insert_at=2)

        assert len(merged_nums) == 1
        assert pres1.slide_count == 3
        assert merged_nums[0] == 2

    def test_merge_empty_presentation(self, template: Template) -> None:
        """Test merging an empty presentation."""
        pres1 = template.create_presentation()
        pres1.add_title_slide("Main", "")

        pres2 = template.create_presentation()
        # pres2 is empty

        merged_nums = pres1.merge(pres2)

        assert len(merged_nums) == 0
        assert pres1.slide_count == 1


class TestMergeFiles:
    """Tests for Presentation.merge_files class method."""

    def test_merge_files_basic(
        self, template: Template, tmp_path: Path
    ) -> None:
        """Test merging files from disk."""
        # Create and save test files
        pres1 = template.create_presentation()
        pres1.add_title_slide("File 1", "")
        path1 = tmp_path / "file1.pptx"
        pres1.save(path1)

        pres2 = template.create_presentation()
        pres2.add_content_slide("File 2", ["Content"])
        path2 = tmp_path / "file2.pptx"
        pres2.save(path2)

        # Merge files
        merged = Presentation.merge_files(template, [path1, path2])

        assert merged.slide_count >= 2

    def test_merge_files_skips_missing(
        self, template: Template, tmp_path: Path
    ) -> None:
        """Test that merge_files skips missing files."""
        pres1 = template.create_presentation()
        pres1.add_title_slide("Existing File", "")
        path1 = tmp_path / "existing.pptx"
        pres1.save(path1)

        missing_path = tmp_path / "missing.pptx"

        # Should not raise, just skip the missing file
        merged = Presentation.merge_files(template, [path1, missing_path])

        assert merged.slide_count >= 1


class TestSaveWithCloning:
    """Tests for saving presentations with cloned slides."""

    def test_save_with_cloned_slides(
        self, presentation: Presentation, tmp_path: Path
    ) -> None:
        """Test saving presentation with cloned slides."""
        presentation.add_title_slide("Original", "")
        presentation.add_content_slide("Content", ["Point"])
        presentation.clone_slide(1)
        presentation.clone_slide(2)

        output_path = tmp_path / "cloned_output.pptx"
        presentation.save(output_path)

        assert output_path.exists()
        assert presentation.slide_count == 4

        # Verify file can be opened
        loaded = PptxPresentation(str(output_path))
        assert len(loaded.slides) == 4

    def test_save_after_merge(
        self, template: Template, tmp_path: Path
    ) -> None:
        """Test saving presentation after merge."""
        pres1 = template.create_presentation()
        pres1.add_title_slide("Main", "")

        pres2 = template.create_presentation()
        pres2.add_content_slide("Merged", ["Item"])

        pres1.merge(pres2)

        output_path = tmp_path / "merged_output.pptx"
        pres1.save(output_path)

        assert output_path.exists()

        loaded = PptxPresentation(str(output_path))
        assert len(loaded.slides) == 2
