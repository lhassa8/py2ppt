"""Tests for shape manipulation features."""

from pathlib import Path

import pytest
from pptx import Presentation as PptxPresentation

from py2ppt import Template, Presentation, ShapeType, ConnectorType


@pytest.fixture
def template(tmp_path: Path) -> Template:
    """Create a template for testing."""
    template_path = tmp_path / "template.pptx"
    pres = PptxPresentation()
    pres.slides.add_slide(pres.slide_layouts[0])
    pres.save(str(template_path))
    return Template(template_path)


@pytest.fixture
def presentation(template: Template) -> Presentation:
    """Create a presentation from template."""
    pres = template.create_presentation()
    pres.add_blank_slide()
    return pres


class TestShapeTypes:
    """Tests for ShapeType enum."""

    def test_shape_type_values(self) -> None:
        """Test ShapeType enum has expected values."""
        assert ShapeType.RECTANGLE.value == "rectangle"
        assert ShapeType.OVAL.value == "oval"
        assert ShapeType.ARROW_RIGHT.value == "arrow_right"

    def test_connector_type_values(self) -> None:
        """Test ConnectorType enum has expected values."""
        assert ConnectorType.STRAIGHT.value == "straight"
        assert ConnectorType.ELBOW.value == "elbow"
        assert ConnectorType.CURVED.value == "curved"


class TestAddTextbox:
    """Tests for add_textbox method."""

    def test_add_basic_textbox(self, presentation: Presentation) -> None:
        """Test adding a basic text box."""
        name = presentation.add_textbox(1, "Hello World", 1, 1, 3, 0.5)

        assert name is not None
        assert isinstance(name, str)

    def test_add_textbox_with_formatting(self, presentation: Presentation) -> None:
        """Test adding a text box with formatting."""
        name = presentation.add_textbox(
            1, "Styled Text",
            left=1, top=2, width=4, height=1,
            font_size=24,
            font_color="#FF0000",
            bold=True,
            italic=True,
            align="center"
        )

        assert name is not None

    def test_add_textbox_invalid_slide(self, presentation: Presentation) -> None:
        """Test adding textbox to invalid slide raises error."""
        from py2ppt.errors import SlideNotFoundError

        with pytest.raises(SlideNotFoundError):
            presentation.add_textbox(99, "Text", 1, 1, 1, 1)


class TestAddShape:
    """Tests for add_shape method."""

    def test_add_rectangle(self, presentation: Presentation) -> None:
        """Test adding a rectangle shape."""
        name = presentation.add_shape(1, "rectangle", 1, 2, 3, 2)

        assert name is not None
        assert isinstance(name, str)

    def test_add_oval(self, presentation: Presentation) -> None:
        """Test adding an oval shape."""
        name = presentation.add_shape(1, "oval", 1, 2, 2, 2)

        assert name is not None

    def test_add_arrow(self, presentation: Presentation) -> None:
        """Test adding an arrow shape."""
        name = presentation.add_shape(1, "arrow_right", 1, 2, 2, 1)

        assert name is not None

    def test_add_shape_with_text(self, presentation: Presentation) -> None:
        """Test adding a shape with text."""
        name = presentation.add_shape(
            1, "rectangle", 1, 2, 3, 2,
            text="Step 1",
            font_size=18,
            font_color="#FFFFFF"
        )

        assert name is not None

    def test_add_shape_with_styling(self, presentation: Presentation) -> None:
        """Test adding a shape with fill and line styling."""
        name = presentation.add_shape(
            1, "rectangle", 1, 2, 3, 2,
            fill_color="#4472C4",
            line_color="#000000",
            line_width=2
        )

        assert name is not None

    def test_add_shape_enum(self, presentation: Presentation) -> None:
        """Test adding shape using ShapeType enum value."""
        name = presentation.add_shape(1, ShapeType.HEXAGON.value, 1, 2, 2, 2)

        assert name is not None

    def test_add_shape_invalid_type(self, presentation: Presentation) -> None:
        """Test adding shape with invalid type raises error."""
        with pytest.raises(ValueError):
            presentation.add_shape(1, "invalid_shape", 1, 1, 1, 1)

    def test_add_multiple_shapes(self, presentation: Presentation) -> None:
        """Test adding multiple shapes."""
        name1 = presentation.add_shape(1, "rectangle", 1, 2, 2, 1)
        name2 = presentation.add_shape(1, "oval", 4, 2, 2, 1)
        name3 = presentation.add_shape(1, "triangle", 7, 2, 2, 1)

        # All shapes should have unique names
        assert len({name1, name2, name3}) == 3


class TestAddConnector:
    """Tests for add_connector method."""

    def test_add_elbow_connector(self, presentation: Presentation) -> None:
        """Test adding an elbow connector."""
        shape1 = presentation.add_shape(1, "rectangle", 1, 2, 2, 1, text="Start")
        shape2 = presentation.add_shape(1, "rectangle", 5, 2, 2, 1, text="End")

        connector = presentation.add_connector(1, shape1, shape2, "elbow")

        assert connector is not None

    def test_add_straight_connector(self, presentation: Presentation) -> None:
        """Test adding a straight connector."""
        shape1 = presentation.add_shape(1, "oval", 1, 2, 1, 1)
        shape2 = presentation.add_shape(1, "oval", 5, 2, 1, 1)

        connector = presentation.add_connector(1, shape1, shape2, "straight")

        assert connector is not None

    def test_add_connector_with_styling(self, presentation: Presentation) -> None:
        """Test adding a connector with styling."""
        shape1 = presentation.add_shape(1, "rectangle", 1, 2, 2, 1)
        shape2 = presentation.add_shape(1, "rectangle", 5, 2, 2, 1)

        connector = presentation.add_connector(
            1, shape1, shape2, "elbow",
            line_color="#FF0000",
            line_width=3
        )

        assert connector is not None

    def test_add_connector_invalid_shape(self, presentation: Presentation) -> None:
        """Test adding connector with invalid shape raises error."""
        from py2ppt.errors import InvalidDataError

        shape1 = presentation.add_shape(1, "rectangle", 1, 2, 2, 1)

        with pytest.raises(InvalidDataError):
            presentation.add_connector(1, shape1, "NonexistentShape", "elbow")


class TestStyleShape:
    """Tests for style_shape method."""

    def test_style_shape_fill(self, presentation: Presentation) -> None:
        """Test styling shape fill color."""
        shape_name = presentation.add_shape(1, "rectangle", 1, 2, 2, 2)

        presentation.style_shape(1, shape_name, fill_color="#FF0000")

        # Should not raise
        info = presentation.get_shape(1, shape_name)
        assert info["name"] == shape_name

    def test_style_shape_line(self, presentation: Presentation) -> None:
        """Test styling shape line color and width."""
        shape_name = presentation.add_shape(1, "rectangle", 1, 2, 2, 2)

        presentation.style_shape(
            1, shape_name,
            line_color="#0000FF",
            line_width=3
        )

        # Should not raise
        info = presentation.get_shape(1, shape_name)
        assert info["name"] == shape_name

    def test_style_shape_invalid_name(self, presentation: Presentation) -> None:
        """Test styling nonexistent shape raises error."""
        from py2ppt.errors import InvalidDataError

        with pytest.raises(InvalidDataError):
            presentation.style_shape(1, "NonexistentShape", fill_color="#FF0000")


class TestGetShape:
    """Tests for get_shape method."""

    def test_get_shape_info(self, presentation: Presentation) -> None:
        """Test getting shape information."""
        shape_name = presentation.add_shape(
            1, "rectangle", 2, 3, 4, 2,
            text="Test Shape"
        )

        info = presentation.get_shape(1, shape_name)

        assert info["name"] == shape_name
        assert "left" in info
        assert "top" in info
        assert "width" in info
        assert "height" in info
        assert info["left_inches"] == pytest.approx(2, rel=0.1)
        assert info["top_inches"] == pytest.approx(3, rel=0.1)
        assert info["text"] == "Test Shape"

    def test_get_shape_invalid_name(self, presentation: Presentation) -> None:
        """Test getting nonexistent shape raises error."""
        from py2ppt.errors import InvalidDataError

        with pytest.raises(InvalidDataError):
            presentation.get_shape(1, "NonexistentShape")


class TestSaveWithShapes:
    """Tests for saving presentations with shapes."""

    def test_save_with_shapes(
        self, presentation: Presentation, tmp_path: Path
    ) -> None:
        """Test saving presentation with shapes."""
        presentation.add_shape(1, "rectangle", 1, 2, 3, 2, fill_color="#4472C4")
        presentation.add_shape(1, "oval", 5, 2, 2, 2, fill_color="#ED7D31")
        presentation.add_textbox(1, "Hello World", 1, 5, 6, 1, font_size=24)

        output_path = tmp_path / "shapes_output.pptx"
        presentation.save(output_path)

        assert output_path.exists()

        # Verify saved file can be opened
        loaded = PptxPresentation(str(output_path))
        assert len(loaded.slides) == 1

        # Check shapes exist
        slide = loaded.slides[0]
        # Should have at least the shapes we added (may have more from template)
        assert len(slide.shapes) >= 3
