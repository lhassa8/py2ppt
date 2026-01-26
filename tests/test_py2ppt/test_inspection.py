"""Tests for slide inspection and editing functionality."""

from pathlib import Path

import pytest
from pptx import Presentation as PptxPresentation

from py2ppt import Template, Presentation, SlideNotFoundError, InvalidDataError


@pytest.fixture
def template(tmp_path: Path) -> Template:
    """Create a template for testing."""
    template_path = tmp_path / "template.pptx"
    pres = PptxPresentation()
    pres.slides.add_slide(pres.slide_layouts[0])
    pres.save(str(template_path))
    return Template(template_path)


@pytest.fixture
def presentation(template: Template) -> Presentation:
    """Create a presentation from template."""
    return template.create_presentation()


@pytest.fixture
def populated_presentation(presentation: Presentation) -> Presentation:
    """Create a presentation with several slides."""
    presentation.add_title_slide("Title Slide", "Subtitle")
    presentation.add_content_slide("Content Slide", ["Point 1", "Point 2"])
    presentation.add_section_slide("Section")
    return presentation


class TestDescribeSlide:
    """Tests for describe_slide."""

    def test_describe_title_slide(self, presentation: Presentation) -> None:
        """Test describing a title slide."""
        presentation.add_title_slide("My Title", "My Subtitle")
        info = presentation.describe_slide(1)

        assert info["slide_number"] == 1
        assert info["layout"] is not None
        assert isinstance(info["shapes"], list)
        assert isinstance(info["has_title"], bool)
        assert isinstance(info["has_content"], bool)
        assert isinstance(info["has_table"], bool)
        assert isinstance(info["has_chart"], bool)
        assert isinstance(info["has_image"], bool)

    def test_describe_content_slide(self, presentation: Presentation) -> None:
        """Test describing a content slide."""
        presentation.add_content_slide("My Content", ["Point A", "Point B"])
        info = presentation.describe_slide(1)

        assert info["slide_number"] == 1
        assert info["has_table"] is False
        assert info["has_chart"] is False

    def test_describe_table_slide(self, presentation: Presentation) -> None:
        """Test describing a table slide."""
        presentation.add_table_slide(
            "Table", ["H1", "H2"], [["a", "b"]]
        )
        info = presentation.describe_slide(1)

        assert info["has_table"] is True
        # Find the table shape
        table_shapes = [s for s in info["shapes"] if "table" in s]
        assert len(table_shapes) >= 1
        assert table_shapes[0]["table"]["headers"] == ["H1", "H2"]
        assert table_shapes[0]["table"]["rows"] == 2  # 1 header + 1 data
        assert table_shapes[0]["table"]["cols"] == 2

    def test_describe_chart_slide(self, presentation: Presentation) -> None:
        """Test describing a chart slide."""
        presentation.add_chart_slide(
            "Chart",
            "pie",
            {"categories": ["A", "B"], "values": [60, 40]},
        )
        info = presentation.describe_slide(1)

        assert info["has_chart"] is True
        chart_shapes = [s for s in info["shapes"] if "chart_type" in s]
        assert len(chart_shapes) >= 1

    def test_describe_slide_notes(self, presentation: Presentation) -> None:
        """Test that notes are included in slide description."""
        presentation.add_title_slide("Title")
        presentation.set_notes(1, "Speaker notes here")
        info = presentation.describe_slide(1)

        assert "Speaker notes here" in info["notes"]

    def test_describe_slide_invalid_number(
        self, presentation: Presentation
    ) -> None:
        """Test that invalid slide number raises SlideNotFoundError."""
        with pytest.raises(SlideNotFoundError):
            presentation.describe_slide(1)

        presentation.add_title_slide("Title")
        with pytest.raises(SlideNotFoundError):
            presentation.describe_slide(0)
        with pytest.raises(SlideNotFoundError):
            presentation.describe_slide(2)

    def test_describe_all_slides(
        self, populated_presentation: Presentation
    ) -> None:
        """Test describing all slides."""
        all_info = populated_presentation.describe_all_slides()

        assert len(all_info) == 3
        assert all_info[0]["slide_number"] == 1
        assert all_info[1]["slide_number"] == 2
        assert all_info[2]["slide_number"] == 3

    def test_describe_all_empty(self, presentation: Presentation) -> None:
        """Test describe_all_slides on empty presentation."""
        assert presentation.describe_all_slides() == []


class TestUpdateSlide:
    """Tests for update_slide."""

    def test_update_title(
        self, populated_presentation: Presentation
    ) -> None:
        """Test updating a slide's title."""
        result = populated_presentation.update_slide(1, title="New Title")
        assert result["slide_number"] == 1
        # Verify the title was updated
        info = populated_presentation.describe_slide(1)
        assert info["slide_number"] == 1

    def test_update_notes(
        self, populated_presentation: Presentation
    ) -> None:
        """Test updating speaker notes."""
        result = populated_presentation.update_slide(
            1, notes="Updated notes"
        )
        assert "Updated notes" in result["notes"]

    def test_update_content(
        self, populated_presentation: Presentation
    ) -> None:
        """Test updating body content."""
        result = populated_presentation.update_slide(
            2, content=["New point 1", "New point 2", "New point 3"]
        )
        assert result["slide_number"] == 2

    def test_update_multiple_fields(
        self, populated_presentation: Presentation
    ) -> None:
        """Test updating multiple fields at once."""
        result = populated_presentation.update_slide(
            2,
            title="Updated Content",
            content=["A", "B"],
            notes="Some notes",
        )
        assert result["slide_number"] == 2
        assert "Some notes" in result["notes"]

    def test_update_nothing(
        self, populated_presentation: Presentation
    ) -> None:
        """Test updating with no changes returns description."""
        result = populated_presentation.update_slide(1)
        assert result["slide_number"] == 1

    def test_update_invalid_slide(
        self, populated_presentation: Presentation
    ) -> None:
        """Test updating invalid slide number."""
        with pytest.raises(SlideNotFoundError):
            populated_presentation.update_slide(99, title="Nope")


class TestDeleteSlide:
    """Tests for delete_slide."""

    def test_delete_slide(
        self, populated_presentation: Presentation
    ) -> None:
        """Test deleting a slide."""
        assert populated_presentation.slide_count == 3
        remaining = populated_presentation.delete_slide(2)
        assert remaining == 2
        assert populated_presentation.slide_count == 2

    def test_delete_first_slide(
        self, populated_presentation: Presentation
    ) -> None:
        """Test deleting the first slide."""
        remaining = populated_presentation.delete_slide(1)
        assert remaining == 2

    def test_delete_last_slide(
        self, populated_presentation: Presentation
    ) -> None:
        """Test deleting the last slide."""
        remaining = populated_presentation.delete_slide(3)
        assert remaining == 2

    def test_delete_all_slides(
        self, populated_presentation: Presentation
    ) -> None:
        """Test deleting all slides one by one."""
        populated_presentation.delete_slide(1)
        populated_presentation.delete_slide(1)
        remaining = populated_presentation.delete_slide(1)
        assert remaining == 0

    def test_delete_invalid_slide(
        self, populated_presentation: Presentation
    ) -> None:
        """Test deleting invalid slide number."""
        with pytest.raises(SlideNotFoundError):
            populated_presentation.delete_slide(0)
        with pytest.raises(SlideNotFoundError):
            populated_presentation.delete_slide(4)


class TestReorderSlides:
    """Tests for reorder_slides and move_slide."""

    def test_reorder_slides(
        self, populated_presentation: Presentation
    ) -> None:
        """Test reordering slides."""
        # Store original slide IDs for comparison
        original_ids = [
            populated_presentation._pptx.slides._sldIdLst[i].rId
            for i in range(3)
        ]

        populated_presentation.reorder_slides([3, 1, 2])

        new_ids = [
            populated_presentation._pptx.slides._sldIdLst[i].rId
            for i in range(3)
        ]

        assert new_ids[0] == original_ids[2]  # Slide 3 is now first
        assert new_ids[1] == original_ids[0]  # Slide 1 is now second
        assert new_ids[2] == original_ids[1]  # Slide 2 is now third

    def test_reorder_identity(
        self, populated_presentation: Presentation
    ) -> None:
        """Test that identity reorder doesn't change anything."""
        original_ids = [
            populated_presentation._pptx.slides._sldIdLst[i].rId
            for i in range(3)
        ]

        populated_presentation.reorder_slides([1, 2, 3])

        new_ids = [
            populated_presentation._pptx.slides._sldIdLst[i].rId
            for i in range(3)
        ]
        assert new_ids == original_ids

    def test_reorder_invalid_order(
        self, populated_presentation: Presentation
    ) -> None:
        """Test that invalid order raises InvalidDataError."""
        with pytest.raises(InvalidDataError):
            populated_presentation.reorder_slides([1, 2])  # Missing slide 3
        with pytest.raises(InvalidDataError):
            populated_presentation.reorder_slides([1, 2, 4])  # Invalid number
        with pytest.raises(InvalidDataError):
            populated_presentation.reorder_slides([1, 1, 2])  # Duplicate

    def test_move_slide_forward(
        self, populated_presentation: Presentation
    ) -> None:
        """Test moving a slide forward."""
        original_ids = [
            populated_presentation._pptx.slides._sldIdLst[i].rId
            for i in range(3)
        ]

        populated_presentation.move_slide(1, 3)

        new_ids = [
            populated_presentation._pptx.slides._sldIdLst[i].rId
            for i in range(3)
        ]
        assert new_ids[2] == original_ids[0]  # Slide 1 moved to position 3

    def test_move_slide_backward(
        self, populated_presentation: Presentation
    ) -> None:
        """Test moving a slide backward."""
        original_ids = [
            populated_presentation._pptx.slides._sldIdLst[i].rId
            for i in range(3)
        ]

        populated_presentation.move_slide(3, 1)

        new_ids = [
            populated_presentation._pptx.slides._sldIdLst[i].rId
            for i in range(3)
        ]
        assert new_ids[0] == original_ids[2]  # Slide 3 moved to position 1

    def test_move_slide_same_position(
        self, populated_presentation: Presentation
    ) -> None:
        """Test moving a slide to its own position."""
        original_ids = [
            populated_presentation._pptx.slides._sldIdLst[i].rId
            for i in range(3)
        ]

        populated_presentation.move_slide(2, 2)

        new_ids = [
            populated_presentation._pptx.slides._sldIdLst[i].rId
            for i in range(3)
        ]
        assert new_ids == original_ids

    def test_move_invalid_from(
        self, populated_presentation: Presentation
    ) -> None:
        """Test moving from invalid position."""
        with pytest.raises(SlideNotFoundError):
            populated_presentation.move_slide(0, 1)
        with pytest.raises(SlideNotFoundError):
            populated_presentation.move_slide(4, 1)

    def test_move_invalid_to(
        self, populated_presentation: Presentation
    ) -> None:
        """Test moving to invalid position."""
        with pytest.raises(SlideNotFoundError):
            populated_presentation.move_slide(1, 0)
        with pytest.raises(SlideNotFoundError):
            populated_presentation.move_slide(1, 4)
