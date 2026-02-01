"""Tests for strategic slide pattern features."""

from pathlib import Path

import pytest
from pptx import Presentation as PptxPresentation

from py2ppt import Template, Presentation


@pytest.fixture
def template(tmp_path: Path) -> Template:
    """Create a template for testing."""
    template_path = tmp_path / "template.pptx"
    pres = PptxPresentation()
    pres.slides.add_slide(pres.slide_layouts[0])
    pres.save(str(template_path))
    return Template(template_path)


@pytest.fixture
def presentation(template: Template) -> Presentation:
    """Create a presentation from template."""
    return template.create_presentation()


class TestSwotSlide:
    """Tests for add_swot_slide method."""

    def test_add_basic_swot(self, presentation: Presentation) -> None:
        """Test adding a basic SWOT slide."""
        slide_num = presentation.add_swot_slide(
            "SWOT Analysis",
            strengths=["Strong brand", "Skilled team"],
            weaknesses=["Limited budget"],
            opportunities=["New markets"],
            threats=["Competition"]
        )

        assert slide_num == 1
        assert presentation.slide_count == 1

    def test_swot_with_empty_quadrants(self, presentation: Presentation) -> None:
        """Test SWOT with some empty quadrants."""
        slide_num = presentation.add_swot_slide(
            "SWOT Analysis",
            strengths=["One strength"],
            weaknesses=[],
            opportunities=["Opportunity 1", "Opportunity 2"],
            threats=[]
        )

        assert slide_num == 1

    def test_swot_with_many_items(self, presentation: Presentation) -> None:
        """Test SWOT with many items in each quadrant."""
        slide_num = presentation.add_swot_slide(
            "Detailed SWOT",
            strengths=["S1", "S2", "S3", "S4", "S5"],
            weaknesses=["W1", "W2", "W3"],
            opportunities=["O1", "O2", "O3", "O4"],
            threats=["T1", "T2"]
        )

        assert slide_num == 1


class TestMatrixSlide:
    """Tests for add_matrix_slide method."""

    def test_add_basic_matrix(self, presentation: Presentation) -> None:
        """Test adding a basic 2x2 matrix."""
        slide_num = presentation.add_matrix_slide(
            "Priority Matrix",
            top_left=["Quick wins"],
            top_right=["Major projects"],
            bottom_left=["Fill-ins"],
            bottom_right=["Low priority"]
        )

        assert slide_num == 1
        assert presentation.slide_count == 1

    def test_matrix_with_axis_labels(self, presentation: Presentation) -> None:
        """Test matrix with axis labels."""
        slide_num = presentation.add_matrix_slide(
            "Impact vs Effort",
            top_left=["High impact, low effort"],
            top_right=["High impact, high effort"],
            bottom_left=["Low impact, low effort"],
            bottom_right=["Low impact, high effort"],
            x_label="Effort",
            y_label="Impact"
        )

        assert slide_num == 1

    def test_matrix_with_quadrant_labels(self, presentation: Presentation) -> None:
        """Test matrix with quadrant labels."""
        slide_num = presentation.add_matrix_slide(
            "Strategic Matrix",
            top_left=["Item 1"],
            top_right=["Item 2"],
            bottom_left=["Item 3"],
            bottom_right=["Item 4"],
            quadrant_labels=("Stars", "Question Marks", "Cash Cows", "Dogs")
        )

        assert slide_num == 1


class TestFunnelSlide:
    """Tests for add_funnel_slide method."""

    def test_add_basic_funnel(self, presentation: Presentation) -> None:
        """Test adding a basic funnel slide."""
        slide_num = presentation.add_funnel_slide(
            "Sales Funnel",
            stages=[
                {"label": "Leads", "value": "1000"},
                {"label": "Qualified", "value": "400"},
                {"label": "Proposals", "value": "100"},
                {"label": "Closed", "value": "25"}
            ]
        )

        assert slide_num == 1
        assert presentation.slide_count == 1

    def test_funnel_with_string_stages(self, presentation: Presentation) -> None:
        """Test funnel with simple string stages."""
        slide_num = presentation.add_funnel_slide(
            "Process Funnel",
            stages=["Awareness", "Interest", "Consideration", "Purchase"]
        )

        assert slide_num == 1

    def test_funnel_with_mixed_stages(self, presentation: Presentation) -> None:
        """Test funnel with mixed dict and string stages."""
        slide_num = presentation.add_funnel_slide(
            "Conversion Funnel",
            stages=[
                {"label": "Visitors", "value": "10K"},
                "Signups",
                {"label": "Active Users", "value": "500"}
            ]
        )

        assert slide_num == 1


class TestPyramidSlide:
    """Tests for add_pyramid_slide method."""

    def test_add_basic_pyramid(self, presentation: Presentation) -> None:
        """Test adding a basic pyramid slide."""
        slide_num = presentation.add_pyramid_slide(
            "Strategic Hierarchy",
            levels=["Vision", "Strategy", "Tactics", "Operations"]
        )

        assert slide_num == 1
        assert presentation.slide_count == 1

    def test_pyramid_with_few_levels(self, presentation: Presentation) -> None:
        """Test pyramid with only 2 levels."""
        slide_num = presentation.add_pyramid_slide(
            "Simple Pyramid",
            levels=["Top", "Bottom"]
        )

        assert slide_num == 1

    def test_pyramid_with_many_levels(self, presentation: Presentation) -> None:
        """Test pyramid with many levels."""
        slide_num = presentation.add_pyramid_slide(
            "Maslow's Hierarchy",
            levels=[
                "Self-Actualization",
                "Esteem",
                "Love/Belonging",
                "Safety",
                "Physiological"
            ]
        )

        assert slide_num == 1


class TestProcessSlide:
    """Tests for add_process_slide method."""

    def test_add_basic_process(self, presentation: Presentation) -> None:
        """Test adding a basic process slide."""
        slide_num = presentation.add_process_slide(
            "Development Process",
            steps=["Plan", "Build", "Test", "Deploy"]
        )

        assert slide_num == 1
        assert presentation.slide_count == 1

    def test_process_with_few_steps(self, presentation: Presentation) -> None:
        """Test process with only 2 steps."""
        slide_num = presentation.add_process_slide(
            "Simple Process",
            steps=["Start", "Finish"]
        )

        assert slide_num == 1

    def test_process_with_many_steps(self, presentation: Presentation) -> None:
        """Test process with many steps."""
        slide_num = presentation.add_process_slide(
            "Extended Process",
            steps=["Initiate", "Plan", "Execute", "Monitor", "Control", "Close"]
        )

        assert slide_num == 1

    def test_process_empty_steps(self, presentation: Presentation) -> None:
        """Test process with empty steps."""
        slide_num = presentation.add_process_slide(
            "Empty Process",
            steps=[]
        )

        assert slide_num == 1


class TestVennSlide:
    """Tests for add_venn_slide method."""

    def test_add_two_circle_venn(self, presentation: Presentation) -> None:
        """Test adding a 2-circle Venn diagram."""
        slide_num = presentation.add_venn_slide(
            "Skills Overlap",
            sets=["Technical", "Creative"],
            intersection_label="Innovation"
        )

        assert slide_num == 1
        assert presentation.slide_count == 1

    def test_add_three_circle_venn(self, presentation: Presentation) -> None:
        """Test adding a 3-circle Venn diagram."""
        slide_num = presentation.add_venn_slide(
            "Role Requirements",
            sets=["Technical", "Business", "Leadership"],
            intersection_label="Ideal Candidate"
        )

        assert slide_num == 1

    def test_venn_without_intersection_label(self, presentation: Presentation) -> None:
        """Test Venn without intersection label."""
        slide_num = presentation.add_venn_slide(
            "Concepts",
            sets=["Set A", "Set B"]
        )

        assert slide_num == 1

    def test_venn_with_single_set(self, presentation: Presentation) -> None:
        """Test Venn with single set (should still work)."""
        slide_num = presentation.add_venn_slide(
            "Single Set",
            sets=["Only One"]
        )

        assert slide_num == 1


class TestSaveWithPatterns:
    """Tests for saving presentations with pattern slides."""

    def test_save_with_all_patterns(
        self, presentation: Presentation, tmp_path: Path
    ) -> None:
        """Test saving presentation with all pattern types."""
        presentation.add_swot_slide(
            "SWOT",
            ["S1"], ["W1"], ["O1"], ["T1"]
        )
        presentation.add_matrix_slide(
            "Matrix",
            ["TL"], ["TR"], ["BL"], ["BR"]
        )
        presentation.add_funnel_slide(
            "Funnel",
            ["Stage 1", "Stage 2", "Stage 3"]
        )
        presentation.add_pyramid_slide(
            "Pyramid",
            ["Top", "Middle", "Bottom"]
        )
        presentation.add_process_slide(
            "Process",
            ["Step 1", "Step 2", "Step 3"]
        )
        presentation.add_venn_slide(
            "Venn",
            ["A", "B", "C"]
        )

        output_path = tmp_path / "patterns_output.pptx"
        presentation.save(output_path)

        assert output_path.exists()
        assert presentation.slide_count == 6

        # Verify saved file can be opened
        loaded = PptxPresentation(str(output_path))
        assert len(loaded.slides) == 6
