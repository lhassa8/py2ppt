"""Tests for export features (PDF, etc.)."""

from pathlib import Path

import pytest
from pptx import Presentation as PptxPresentation

from py2ppt import Template, Presentation, is_pdf_export_available


@pytest.fixture
def template(tmp_path: Path) -> Template:
    """Create a template for testing."""
    template_path = tmp_path / "template.pptx"
    pres = PptxPresentation()
    pres.slides.add_slide(pres.slide_layouts[0])
    pres.save(str(template_path))
    return Template(template_path)


@pytest.fixture
def presentation(template: Template) -> Presentation:
    """Create a presentation from template."""
    return template.create_presentation()


class TestIsPdfExportAvailable:
    """Tests for is_pdf_export_available function."""

    def test_returns_dict(self) -> None:
        """Test that function returns a dict."""
        available = is_pdf_export_available()

        assert isinstance(available, dict)
        assert "libreoffice" in available
        assert "unoconv" in available

    def test_values_are_bool(self) -> None:
        """Test that values are boolean."""
        available = is_pdf_export_available()

        assert isinstance(available["libreoffice"], bool)
        assert isinstance(available["unoconv"], bool)


class TestSavePdf:
    """Tests for save_pdf method."""

    def test_save_pdf_without_libreoffice(
        self, presentation: Presentation, tmp_path: Path
    ) -> None:
        """Test PDF export behavior when LibreOffice not available."""
        from py2ppt.export import ExportError

        available = is_pdf_export_available()

        presentation.add_title_slide("PDF Test", "")
        output_path = tmp_path / "output.pdf"

        if not available["libreoffice"]:
            # Should raise ExportError when LibreOffice not available
            with pytest.raises(ExportError):
                presentation.save_pdf(output_path)
        else:
            # If LibreOffice is available, should succeed
            presentation.save_pdf(output_path)
            assert output_path.exists()

    @pytest.mark.skipif(
        not is_pdf_export_available()["libreoffice"],
        reason="LibreOffice not installed"
    )
    def test_save_pdf_basic(
        self, presentation: Presentation, tmp_path: Path
    ) -> None:
        """Test basic PDF export (requires LibreOffice)."""
        presentation.add_title_slide("PDF Export Test", "Subtitle")
        presentation.add_content_slide("Content", ["Point 1", "Point 2"])

        output_path = tmp_path / "test_output.pdf"
        presentation.save_pdf(output_path)

        assert output_path.exists()
        assert output_path.stat().st_size > 0

    @pytest.mark.skipif(
        not is_pdf_export_available()["libreoffice"],
        reason="LibreOffice not installed"
    )
    def test_save_pdf_multiple_slides(
        self, presentation: Presentation, tmp_path: Path
    ) -> None:
        """Test PDF export with multiple slides."""
        for i in range(5):
            presentation.add_content_slide(f"Slide {i+1}", [f"Content {i+1}"])

        output_path = tmp_path / "multi_slide.pdf"
        presentation.save_pdf(output_path)

        assert output_path.exists()

    def test_save_pdf_invalid_engine(
        self, presentation: Presentation, tmp_path: Path
    ) -> None:
        """Test PDF export with invalid engine raises error."""
        from py2ppt.export import ExportError

        presentation.add_title_slide("Test", "")
        output_path = tmp_path / "output.pdf"

        with pytest.raises(ExportError):
            presentation.save_pdf(output_path, engine="invalid_engine")


class TestExportError:
    """Tests for ExportError exception."""

    def test_export_error_message(self) -> None:
        """Test ExportError has correct structure."""
        from py2ppt.export import ExportError

        error = ExportError(
            "Test error message",
            suggestion="Try this fix",
            code="TEST_ERROR",
        )

        assert str(error) == "Test error message"
        assert error.suggestion == "Try this fix"
        assert error.code == "TEST_ERROR"

    def test_export_error_to_dict(self) -> None:
        """Test ExportError can be converted to dict."""
        from py2ppt.export import ExportError

        error = ExportError(
            "Error message",
            suggestion="Suggestion",
            code="CODE",
        )

        d = error.to_dict()

        assert d["message"] == "Error message"
        assert d["suggestion"] == "Suggestion"
        assert d["code"] == "CODE"
