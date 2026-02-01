"""Tests for presentation diff/comparison features."""

from pathlib import Path

import pytest
from pptx import Presentation as PptxPresentation

from py2ppt import Template, Presentation, diff_presentations


@pytest.fixture
def template(tmp_path: Path) -> Template:
    """Create a template for testing."""
    template_path = tmp_path / "template.pptx"
    pres = PptxPresentation()
    pres.slides.add_slide(pres.slide_layouts[0])
    pres.save(str(template_path))
    return Template(template_path)


class TestDiffPresentations:
    """Tests for diff_presentations function."""

    def test_diff_identical_presentations(self, template: Template) -> None:
        """Test diffing identical presentations."""
        pres1 = template.create_presentation()
        pres1.add_title_slide("Title", "Subtitle")

        pres2 = template.create_presentation()
        pres2.add_title_slide("Title", "Subtitle")

        diff = diff_presentations(pres1, pres2)

        assert diff["slides_added"] == []
        assert diff["slides_removed"] == []
        assert diff["slides_modified"] == []

    def test_diff_added_slides(self, template: Template) -> None:
        """Test detecting added slides."""
        pres1 = template.create_presentation()
        pres1.add_title_slide("Title", "")

        pres2 = template.create_presentation()
        pres2.add_title_slide("Title", "")
        pres2.add_content_slide("New Slide", ["Content"])

        diff = diff_presentations(pres1, pres2)

        assert 2 in diff["slides_added"]
        assert len(diff["slides_removed"]) == 0

    def test_diff_removed_slides(self, template: Template) -> None:
        """Test detecting removed slides."""
        pres1 = template.create_presentation()
        pres1.add_title_slide("Title", "")
        pres1.add_content_slide("Removed Slide", ["Content"])

        pres2 = template.create_presentation()
        pres2.add_title_slide("Title", "")

        diff = diff_presentations(pres1, pres2)

        assert 2 in diff["slides_removed"]

    def test_diff_modified_title(self, template: Template) -> None:
        """Test detecting title modification."""
        pres1 = template.create_presentation()
        pres1.add_content_slide("Original Title", ["Content"])

        pres2 = template.create_presentation()
        pres2.add_content_slide("Modified Title", ["Content"])

        diff = diff_presentations(pres1, pres2)

        assert len(diff["slides_modified"]) > 0
        modified = diff["slides_modified"][0]
        assert any("title" in c.lower() for c in modified["changes"])

    def test_diff_modified_content(self, template: Template) -> None:
        """Test detecting content modification."""
        pres1 = template.create_presentation()
        pres1.add_content_slide("Title", ["Point 1"])

        pres2 = template.create_presentation()
        pres2.add_content_slide("Title", ["Point 1", "Point 2"])

        diff = diff_presentations(pres1, pres2)

        assert len(diff["slides_modified"]) > 0
        modified = diff["slides_modified"][0]
        assert any("content" in c.lower() for c in modified["changes"])

    def test_diff_text_format(self, template: Template) -> None:
        """Test diff with text format output."""
        pres1 = template.create_presentation()
        pres1.add_title_slide("Title", "")

        pres2 = template.create_presentation()
        pres2.add_title_slide("Title", "")
        pres2.add_content_slide("New", [])

        diff = diff_presentations(pres1, pres2, format="text")

        assert isinstance(diff, str)
        assert "Presentation Diff" in diff
        assert "added" in diff.lower()

    def test_diff_empty_presentations(self, template: Template) -> None:
        """Test diffing empty presentations."""
        pres1 = template.create_presentation()
        pres2 = template.create_presentation()

        diff = diff_presentations(pres1, pres2)

        assert diff["slides_added"] == []
        assert diff["slides_removed"] == []
        assert diff["slides_modified"] == []

    def test_diff_summary(self, template: Template) -> None:
        """Test that diff includes a summary."""
        pres1 = template.create_presentation()
        pres1.add_title_slide("Title", "")

        pres2 = template.create_presentation()
        pres2.add_title_slide("Title", "")
        pres2.add_content_slide("Added", [])

        diff = diff_presentations(pres1, pres2)

        assert "summary" in diff
        assert "added" in diff["summary"].lower()


class TestDescribeMaster:
    """Tests for describe_master method."""

    def test_describe_master_basic(self, template: Template) -> None:
        """Test basic master description."""
        pres = template.create_presentation()

        master = pres.describe_master()

        assert "name" in master
        assert "layout_count" in master
        assert "colors" in master
        assert "fonts" in master

    def test_describe_master_has_layouts(self, template: Template) -> None:
        """Test that master reports layout count."""
        pres = template.create_presentation()

        master = pres.describe_master()

        assert master["layout_count"] >= 1


class TestDescribeLayouts:
    """Tests for describe_layouts method."""

    def test_describe_layouts_basic(self, template: Template) -> None:
        """Test basic layout description."""
        pres = template.create_presentation()

        layouts = pres.describe_layouts()

        assert isinstance(layouts, list)
        assert len(layouts) >= 1

    def test_describe_layouts_has_info(self, template: Template) -> None:
        """Test that layouts include expected info."""
        pres = template.create_presentation()

        layouts = pres.describe_layouts()

        for layout in layouts:
            assert "name" in layout
            assert "index" in layout


class TestGetLayout:
    """Tests for get_layout method."""

    def test_get_layout_by_index(self, template: Template) -> None:
        """Test getting layout by index."""
        pres = template.create_presentation()

        layout = pres.get_layout(0)

        assert layout is not None
        assert "name" in layout

    def test_get_layout_by_name(self, template: Template) -> None:
        """Test getting layout by name."""
        pres = template.create_presentation()
        layouts = pres.describe_layouts()

        if layouts:
            first_name = layouts[0]["name"]
            layout = pres.get_layout(first_name)

            assert layout is not None

    def test_get_layout_not_found(self, template: Template) -> None:
        """Test getting nonexistent layout returns None."""
        pres = template.create_presentation()

        layout = pres.get_layout("Nonexistent Layout Name")

        assert layout is None


class TestThemeModification:
    """Tests for theme modification features."""

    def test_set_theme_color(self, template: Template) -> None:
        """Test setting a theme color."""
        pres = template.create_presentation()

        # Should not raise
        pres.set_theme_color("accent1", "#FF6600")

    def test_save_as_template(
        self, template: Template, tmp_path: Path
    ) -> None:
        """Test saving as a template."""
        pres = template.create_presentation()
        pres.add_title_slide("Template Test", "")

        output_path = tmp_path / "new_template.pptx"
        pres.save_as_template(output_path)

        assert output_path.exists()

        # Verify it can be loaded as a template
        new_template = Template(output_path)
        assert new_template is not None
