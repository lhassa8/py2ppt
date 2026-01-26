"""Tests for chart slide functionality."""

from pathlib import Path

import pytest
from pptx import Presentation as PptxPresentation

from py2ppt import Template, Presentation, InvalidDataError


@pytest.fixture
def template(tmp_path: Path) -> Template:
    """Create a template for testing."""
    template_path = tmp_path / "template.pptx"
    pres = PptxPresentation()
    pres.slides.add_slide(pres.slide_layouts[0])
    pres.save(str(template_path))
    return Template(template_path)


@pytest.fixture
def presentation(template: Template) -> Presentation:
    """Create a presentation from template."""
    return template.create_presentation()


class TestChartSlide:
    """Tests for add_chart_slide."""

    def test_column_chart(self, presentation: Presentation) -> None:
        """Test creating a column chart."""
        slide_num = presentation.add_chart_slide(
            "Revenue",
            "column",
            {
                "categories": ["Q1", "Q2", "Q3"],
                "series": [{"name": "2024", "values": [10, 20, 30]}],
            },
        )
        assert slide_num == 1
        assert presentation.slide_count == 1

    def test_bar_chart(self, presentation: Presentation) -> None:
        """Test creating a bar chart."""
        slide_num = presentation.add_chart_slide(
            "Comparison",
            "bar",
            {
                "categories": ["A", "B", "C"],
                "series": [{"name": "Data", "values": [5, 10, 15]}],
            },
        )
        assert slide_num == 1

    def test_line_chart(self, presentation: Presentation) -> None:
        """Test creating a line chart."""
        slide_num = presentation.add_chart_slide(
            "Trend",
            "line",
            {
                "categories": ["Jan", "Feb", "Mar"],
                "series": [{"name": "Users", "values": [100, 150, 200]}],
            },
        )
        assert slide_num == 1

    def test_pie_chart(self, presentation: Presentation) -> None:
        """Test creating a pie chart."""
        slide_num = presentation.add_chart_slide(
            "Distribution",
            "pie",
            {
                "categories": ["Desktop", "Mobile", "Tablet"],
                "values": [60, 30, 10],
            },
        )
        assert slide_num == 1

    def test_doughnut_chart(self, presentation: Presentation) -> None:
        """Test creating a doughnut chart."""
        slide_num = presentation.add_chart_slide(
            "Breakdown",
            "doughnut",
            {
                "categories": ["Yes", "No"],
                "values": [70, 30],
            },
        )
        assert slide_num == 1

    def test_multi_series_chart(self, presentation: Presentation) -> None:
        """Test multi-series chart gets legend."""
        presentation.add_chart_slide(
            "Multi-Series",
            "column",
            {
                "categories": ["Q1", "Q2"],
                "series": [
                    {"name": "2023", "values": [10, 20]},
                    {"name": "2024", "values": [15, 25]},
                ],
            },
        )
        slide = presentation._pptx.slides[0]
        for shape in slide.shapes:
            if shape.has_chart:
                assert shape.chart.has_legend is True

    def test_single_series_no_legend(self, presentation: Presentation) -> None:
        """Test single-series chart does not get legend."""
        presentation.add_chart_slide(
            "Single",
            "column",
            {
                "categories": ["Q1", "Q2"],
                "series": [{"name": "Revenue", "values": [10, 20]}],
            },
        )
        slide = presentation._pptx.slides[0]
        for shape in slide.shapes:
            if shape.has_chart:
                assert shape.chart.has_legend is False

    def test_invalid_chart_type(self, presentation: Presentation) -> None:
        """Test that invalid chart type raises InvalidDataError."""
        with pytest.raises(InvalidDataError) as exc_info:
            presentation.add_chart_slide(
                "Bad",
                "radar",
                {"categories": ["A"], "series": [{"name": "X", "values": [1]}]},
            )
        assert "INVALID_CHART_TYPE" in exc_info.value.code

    def test_missing_categories(self, presentation: Presentation) -> None:
        """Test that missing categories raises InvalidDataError."""
        with pytest.raises(InvalidDataError) as exc_info:
            presentation.add_chart_slide(
                "Bad",
                "column",
                {"series": [{"name": "X", "values": [1]}]},
            )
        assert "MISSING_CHART_DATA" in exc_info.value.code

    def test_pie_missing_values(self, presentation: Presentation) -> None:
        """Test that pie chart without values raises InvalidDataError."""
        with pytest.raises(InvalidDataError) as exc_info:
            presentation.add_chart_slide(
                "Bad Pie",
                "pie",
                {"categories": ["A", "B"]},
            )
        assert "MISSING_CHART_VALUES" in exc_info.value.code

    def test_bar_missing_series(self, presentation: Presentation) -> None:
        """Test that bar chart without series raises InvalidDataError."""
        with pytest.raises(InvalidDataError) as exc_info:
            presentation.add_chart_slide(
                "Bad Bar",
                "bar",
                {"categories": ["A", "B"]},
            )
        assert "MISSING_CHART_SERIES" in exc_info.value.code

    def test_chart_via_add_slide(self, presentation: Presentation) -> None:
        """Test creating a chart via the add_slide auto-router."""
        slide_num = presentation.add_slide(
            content_type="chart",
            title="Auto Chart",
            chart_type="pie",
            data={"categories": ["X", "Y"], "values": [50, 50]},
        )
        assert slide_num == 1

    def test_chart_save_and_reload(
        self, presentation: Presentation, tmp_path: Path
    ) -> None:
        """Test that chart slides survive save/reload."""
        presentation.add_chart_slide(
            "Persist",
            "column",
            {
                "categories": ["A", "B"],
                "series": [{"name": "S", "values": [1, 2]}],
            },
        )
        output = tmp_path / "chart_output.pptx"
        presentation.save(output)

        loaded = PptxPresentation(str(output))
        assert len(loaded.slides) == 1
        found_chart = False
        for shape in loaded.slides[0].shapes:
            if shape.has_chart:
                found_chart = True
        assert found_chart

    def test_chart_case_insensitive(self, presentation: Presentation) -> None:
        """Test that chart_type is case-insensitive."""
        slide_num = presentation.add_chart_slide(
            "Case Test",
            "PIE",
            {"categories": ["A", "B"], "values": [50, 50]},
        )
        assert slide_num == 1
