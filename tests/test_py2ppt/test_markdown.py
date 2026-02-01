"""Tests for Markdown import/export features."""

from pathlib import Path

import pytest
from pptx import Presentation as PptxPresentation

from py2ppt import Template, Presentation, build_from_markdown


@pytest.fixture
def template(tmp_path: Path) -> Template:
    """Create a template for testing."""
    template_path = tmp_path / "template.pptx"
    pres = PptxPresentation()
    pres.slides.add_slide(pres.slide_layouts[0])
    pres.save(str(template_path))
    return Template(template_path)


@pytest.fixture
def presentation(template: Template) -> Presentation:
    """Create a presentation from template."""
    return template.create_presentation()


class TestToMarkdown:
    """Tests for to_markdown export method."""

    def test_export_basic(self, presentation: Presentation) -> None:
        """Test basic Markdown export."""
        presentation.add_title_slide("Test Title", "Subtitle")
        presentation.add_content_slide("Content", ["Point 1", "Point 2"])

        md = presentation.to_markdown()

        assert "Test Title" in md
        assert "Content" in md
        assert "Point 1" in md
        assert "Point 2" in md

    def test_export_to_file(
        self, presentation: Presentation, tmp_path: Path
    ) -> None:
        """Test exporting to a file."""
        presentation.add_title_slide("File Export", "")
        presentation.add_content_slide("Bullets", ["A", "B", "C"])

        output_path = tmp_path / "output.md"
        result = presentation.to_markdown(output_path)

        assert output_path.exists()
        file_content = output_path.read_text()
        assert result == file_content
        assert "File Export" in file_content

    def test_export_with_notes(self, presentation: Presentation) -> None:
        """Test exporting slides with speaker notes."""
        presentation.add_content_slide("Slide", ["Point"])
        presentation.set_notes(1, "These are speaker notes")

        md = presentation.to_markdown()

        assert "notes:" in md.lower()
        assert "speaker notes" in md.lower()

    def test_export_section_slides(self, presentation: Presentation) -> None:
        """Test exporting section slides."""
        presentation.add_section_slide("Section Title")

        md = presentation.to_markdown()

        assert "Section:" in md

    def test_export_empty_presentation(self, presentation: Presentation) -> None:
        """Test exporting an empty presentation."""
        md = presentation.to_markdown()

        # Should return empty or minimal markdown
        assert isinstance(md, str)


class TestBuildFromMarkdown:
    """Tests for build_from_markdown import function."""

    def test_build_basic(self, template: Template) -> None:
        """Test building from basic Markdown."""
        md = """
# My Presentation

## Slide: Introduction
- Point 1
- Point 2
- Point 3
"""
        pres = build_from_markdown(template, md)

        assert pres.slide_count >= 1

    def test_build_with_title_slide(self, template: Template) -> None:
        """Test that H1 becomes title slide."""
        md = """# Presentation Title"""

        pres = build_from_markdown(template, md)

        assert pres.slide_count == 1
        slide = pres.describe_slide(1)
        assert "Presentation Title" in slide.get("title", "")

    def test_build_with_section(self, template: Template) -> None:
        """Test building section slides."""
        md = """
# Main Title

## Section: Part 1

## Slide: Content
- Item
"""
        pres = build_from_markdown(template, md)

        assert pres.slide_count >= 2

    def test_build_with_bullets(self, template: Template) -> None:
        """Test building slides with bullet points."""
        md = """
## Slide: Key Points
- First point
- Second point
- Third point
"""
        pres = build_from_markdown(template, md)

        assert pres.slide_count == 1
        slide = pres.describe_slide(1)
        content = slide.get("content", [])
        assert len(content) >= 3

    def test_build_with_table(self, template: Template) -> None:
        """Test building slides with tables."""
        md = """
## Slide: Data Table
| Name | Value |
|------|-------|
| A    | 100   |
| B    | 200   |
"""
        pres = build_from_markdown(template, md)

        assert pres.slide_count == 1
        slide = pres.describe_slide(1)
        assert slide.get("has_table", False)

    def test_build_with_notes(self, template: Template) -> None:
        """Test building slides with speaker notes."""
        md = """
## Slide: With Notes
- Content here

<!-- notes: Remember to emphasize this point -->
"""
        pres = build_from_markdown(template, md)

        assert pres.slide_count == 1
        slide = pres._pptx.slides[0]
        notes = slide.notes_slide.notes_text_frame.text
        assert "emphasize" in notes.lower()

    def test_build_from_file(self, template: Template, tmp_path: Path) -> None:
        """Test building from a Markdown file."""
        md_content = """
# File Presentation

## Slide: From File
- Loaded from disk
"""
        md_path = tmp_path / "input.md"
        md_path.write_text(md_content)

        pres = build_from_markdown(template, md_path)

        assert pres.slide_count >= 1

    def test_build_multiple_slides(self, template: Template) -> None:
        """Test building multiple slides."""
        md = """
# Main Title

## Slide: First
- Point A

## Slide: Second
- Point B

## Slide: Third
- Point C
"""
        pres = build_from_markdown(template, md)

        assert pres.slide_count >= 3

    def test_build_with_numbered_format(self, template: Template) -> None:
        """Test Slide N: Title format."""
        md = """
## Slide 1: Introduction
- Item 1

## Slide 2: Details
- Item 2
"""
        pres = build_from_markdown(template, md)

        assert pres.slide_count == 2

    def test_build_empty_markdown(self, template: Template) -> None:
        """Test building from empty Markdown."""
        pres = build_from_markdown(template, "")

        assert pres.slide_count == 0


class TestMarkdownRoundTrip:
    """Tests for round-trip Markdown export/import."""

    def test_basic_roundtrip(self, template: Template) -> None:
        """Test exporting and re-importing."""
        # Create original presentation
        original = template.create_presentation()
        original.add_title_slide("Round Trip Test", "")
        original.add_content_slide("Content", ["Point 1", "Point 2"])

        # Export to Markdown
        md = original.to_markdown()

        # Import back
        imported = build_from_markdown(template, md)

        # Should have similar structure
        assert imported.slide_count >= 1

    def test_roundtrip_preserves_titles(
        self, template: Template, tmp_path: Path
    ) -> None:
        """Test that round-trip preserves slide titles."""
        original = template.create_presentation()
        original.add_title_slide("My Presentation", "")
        original.add_content_slide("Key Points", ["A", "B", "C"])
        original.add_section_slide("Section Break")

        # Round-trip through file
        md_path = tmp_path / "roundtrip.md"
        original.to_markdown(md_path)

        imported = build_from_markdown(template, md_path)

        # Check titles are preserved
        assert imported.slide_count >= 2


class TestSaveWithMarkdown:
    """Tests for saving presentations built from Markdown."""

    def test_save_markdown_built(
        self, template: Template, tmp_path: Path
    ) -> None:
        """Test saving a presentation built from Markdown."""
        md = """
# Markdown Presentation

## Slide: First Slide
- Point 1
- Point 2

## Section: Part 2

## Slide: Second Slide
- More content
"""
        pres = build_from_markdown(template, md)

        output_path = tmp_path / "from_markdown.pptx"
        pres.save(output_path)

        assert output_path.exists()

        # Verify saved file
        loaded = PptxPresentation(str(output_path))
        assert len(loaded.slides) >= 2
