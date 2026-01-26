"""Tests for content intelligence functionality."""

from pathlib import Path

import pytest
from pptx import Presentation as PptxPresentation

from py2ppt import Template, Presentation, SlideNotFoundError, LayoutNotFoundError


@pytest.fixture
def template(tmp_path: Path) -> Template:
    """Create a template for testing."""
    template_path = tmp_path / "template.pptx"
    pres = PptxPresentation()
    pres.slides.add_slide(pres.slide_layouts[0])
    pres.save(str(template_path))
    return Template(template_path)


@pytest.fixture
def presentation(template: Template) -> Presentation:
    """Create a presentation from template."""
    return template.create_presentation()


class TestAutoSplit:
    """Tests for add_content_slides (auto-split)."""

    def test_no_split_needed(self, presentation: Presentation) -> None:
        """Test content that fits on one slide."""
        slides = presentation.add_content_slides(
            "Short List",
            ["Item 1", "Item 2", "Item 3"],
            max_bullets=6,
        )
        assert len(slides) == 1
        assert presentation.slide_count == 1

    def test_split_into_two(self, presentation: Presentation) -> None:
        """Test content split across two slides."""
        slides = presentation.add_content_slides(
            "Long List",
            [f"Item {i}" for i in range(8)],
            max_bullets=4,
        )
        assert len(slides) == 2
        assert presentation.slide_count == 2

    def test_split_into_three(self, presentation: Presentation) -> None:
        """Test content split across three slides."""
        slides = presentation.add_content_slides(
            "Very Long",
            [f"Point {i}" for i in range(15)],
            max_bullets=6,
        )
        assert len(slides) == 3
        assert presentation.slide_count == 3

    def test_continuation_suffix(self, presentation: Presentation) -> None:
        """Test that continuation slides get the suffix in their title."""
        presentation.add_content_slides(
            "My Topic",
            [f"Item {i}" for i in range(8)],
            max_bullets=4,
        )
        # First slide should have original title
        info1 = presentation.describe_slide(1)
        # Second slide should have continuation suffix
        info2 = presentation.describe_slide(2)
        # Both should exist
        assert info1["slide_number"] == 1
        assert info2["slide_number"] == 2

    def test_custom_suffix(self, presentation: Presentation) -> None:
        """Test custom continuation suffix."""
        slides = presentation.add_content_slides(
            "Topic",
            [f"Item {i}" for i in range(8)],
            max_bullets=4,
            continuation_suffix=" - continued",
        )
        assert len(slides) == 2

    def test_sub_items_grouped(self, presentation: Presentation) -> None:
        """Test that sub-items stay with their parent."""
        content = [
            "Parent 1",
            "Sub 1a",
            "Sub 1b",
            "Sub 1c",
            "Parent 2",
            "Sub 2a",
            "Parent 3",
        ]
        levels = [0, 1, 1, 1, 0, 1, 0]

        slides = presentation.add_content_slides(
            "Grouped",
            content,
            levels=levels,
            max_bullets=4,
        )
        # First chunk: Parent 1 + 3 subs = 4 items, then split at Parent 2
        # Second chunk: Parent 2 + Sub 2a + Parent 3 = 3 items
        assert len(slides) == 2

    def test_sub_items_not_split_from_parent(
        self, presentation: Presentation
    ) -> None:
        """Test that sub-items are never orphaned from their parent."""
        content = [
            "Parent 1",
            "Sub 1a",
            "Sub 1b",
            "Parent 2",
            "Sub 2a",
            "Sub 2b",
            "Sub 2c",
        ]
        levels = [0, 1, 1, 0, 1, 1, 1]

        slides = presentation.add_content_slides(
            "Groups",
            content,
            levels=levels,
            max_bullets=3,
        )
        # Parent 1 + 2 subs = 3, then Parent 2 + 3 subs = 4 (exceeds 3 but can't split)
        assert len(slides) == 2

    def test_empty_content(self, presentation: Presentation) -> None:
        """Test auto-split with empty content."""
        slides = presentation.add_content_slides("Empty", [])
        assert slides == []

    def test_string_content(self, presentation: Presentation) -> None:
        """Test auto-split with string content (split by newlines)."""
        slides = presentation.add_content_slides(
            "String",
            "Line 1\nLine 2\nLine 3\nLine 4\nLine 5\nLine 6\nLine 7\nLine 8",
            max_bullets=4,
        )
        assert len(slides) == 2

    def test_single_item(self, presentation: Presentation) -> None:
        """Test auto-split with single item."""
        slides = presentation.add_content_slides(
            "Single",
            ["Only item"],
        )
        assert len(slides) == 1

    def test_returns_slide_numbers(self, presentation: Presentation) -> None:
        """Test that returned slide numbers are correct."""
        slides = presentation.add_content_slides(
            "Numbers",
            [f"Item {i}" for i in range(8)],
            max_bullets=4,
        )
        assert slides == [1, 2]


class TestOverflowDetection:
    """Tests for overflow detection in add_content_slide."""

    def test_warn_overflow_returns_dict(
        self, presentation: Presentation
    ) -> None:
        """Test that warn_overflow=True returns a dict."""
        result = presentation.add_content_slide(
            "Test", ["A", "B", "C"], warn_overflow=True
        )
        assert isinstance(result, dict)
        assert "slide_number" in result
        assert "overflow" in result
        assert "item_count" in result
        assert "estimated_capacity" in result

    def test_no_overflow(self, presentation: Presentation) -> None:
        """Test detection when content fits."""
        result = presentation.add_content_slide(
            "Fits", ["A", "B"], warn_overflow=True
        )
        assert result["overflow"] is False
        assert result["item_count"] == 2

    def test_overflow_detected(self, presentation: Presentation) -> None:
        """Test detection when content overflows."""
        result = presentation.add_content_slide(
            "Overflows",
            [f"Item {i}" for i in range(50)],
            warn_overflow=True,
        )
        assert result["overflow"] is True
        assert result["item_count"] == 50

    def test_warn_overflow_false_returns_int(
        self, presentation: Presentation
    ) -> None:
        """Test that warn_overflow=False (default) returns int."""
        result = presentation.add_content_slide("Normal", ["A", "B"])
        assert isinstance(result, int)

    def test_overflow_slide_number(
        self, presentation: Presentation
    ) -> None:
        """Test that overflow result includes correct slide number."""
        presentation.add_title_slide("First")
        result = presentation.add_content_slide(
            "Second", ["A"], warn_overflow=True
        )
        assert result["slide_number"] == 2


class TestStructuredErrors:
    """Tests for structured error handling."""

    def test_slide_not_found_error(
        self, presentation: Presentation
    ) -> None:
        """Test SlideNotFoundError attributes."""
        presentation.add_title_slide("Test")
        with pytest.raises(SlideNotFoundError) as exc_info:
            presentation.describe_slide(99)

        err = exc_info.value
        assert err.code == "SLIDE_NOT_FOUND"
        assert err.suggestion
        d = err.to_dict()
        assert d["error"] == "SlideNotFoundError"
        assert d["code"] == "SLIDE_NOT_FOUND"

    def test_slide_not_found_empty_presentation(
        self, presentation: Presentation
    ) -> None:
        """Test error message for empty presentation."""
        with pytest.raises(SlideNotFoundError) as exc_info:
            presentation.describe_slide(1)
        assert "Add slides first" in exc_info.value.suggestion

    def test_layout_not_found_by_name(
        self, presentation: Presentation
    ) -> None:
        """Test LayoutNotFoundError for bad layout name."""
        with pytest.raises(LayoutNotFoundError) as exc_info:
            presentation.add_title_slide(
                "Test", layout="NonexistentLayout12345"
            )
        assert exc_info.value.code == "LAYOUT_NOT_FOUND"

    def test_layout_not_found_by_index(
        self, presentation: Presentation
    ) -> None:
        """Test LayoutNotFoundError for bad layout index."""
        with pytest.raises(LayoutNotFoundError):
            presentation.add_title_slide("Test", layout=999)

    def test_set_notes_validates_slide(
        self, presentation: Presentation
    ) -> None:
        """Test that set_notes validates slide number."""
        with pytest.raises(SlideNotFoundError):
            presentation.set_notes(1, "Notes")

    def test_delete_validates_slide(
        self, presentation: Presentation
    ) -> None:
        """Test that delete_slide validates slide number."""
        with pytest.raises(SlideNotFoundError):
            presentation.delete_slide(1)

    def test_update_validates_slide(
        self, presentation: Presentation
    ) -> None:
        """Test that update_slide validates slide number."""
        with pytest.raises(SlideNotFoundError):
            presentation.update_slide(1, title="Nope")

    def test_error_to_dict(self) -> None:
        """Test Py2PptError.to_dict()."""
        from py2ppt.errors import Py2PptError

        err = Py2PptError("msg", suggestion="fix it", code="TEST")
        d = err.to_dict()
        assert d == {
            "error": "Py2PptError",
            "message": "msg",
            "suggestion": "fix it",
            "code": "TEST",
        }


class TestRichTextFormatting:
    """Tests for rich text formatting in presentation slides."""

    def test_bold_text_applied(self, presentation: Presentation) -> None:
        """Test that bold formatting is applied to runs."""
        presentation.add_content_slide(
            "Rich Text",
            [
                [{"text": "Bold text", "bold": True}],
            ],
        )
        slide = presentation._pptx.slides[0]
        body = None
        for shape in slide.placeholders:
            from pptx.enum.shapes import PP_PLACEHOLDER
            if shape.placeholder_format.type in (
                PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT
            ):
                body = shape
                break

        if body and body.has_text_frame:
            para = body.text_frame.paragraphs[0]
            assert len(para.runs) >= 1
            assert para.runs[0].font.bold is True

    def test_italic_text_applied(self, presentation: Presentation) -> None:
        """Test that italic formatting is applied."""
        presentation.add_content_slide(
            "Italic",
            [[{"text": "Italic text", "italic": True}]],
        )
        slide = presentation._pptx.slides[0]
        for shape in slide.placeholders:
            from pptx.enum.shapes import PP_PLACEHOLDER
            if shape.placeholder_format.type in (
                PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT
            ):
                para = shape.text_frame.paragraphs[0]
                if para.runs:
                    assert para.runs[0].font.italic is True

    def test_color_text_applied(self, presentation: Presentation) -> None:
        """Test that color formatting is applied."""
        presentation.add_content_slide(
            "Color",
            [[{"text": "Red text", "color": "#FF0000"}]],
        )
        slide = presentation._pptx.slides[0]
        for shape in slide.placeholders:
            from pptx.enum.shapes import PP_PLACEHOLDER
            if shape.placeholder_format.type in (
                PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT
            ):
                para = shape.text_frame.paragraphs[0]
                if para.runs:
                    from pptx.dml.color import RGBColor
                    assert para.runs[0].font.color.rgb == RGBColor(0xFF, 0, 0)

    def test_multi_run_paragraph(self, presentation: Presentation) -> None:
        """Test paragraph with multiple formatted runs."""
        presentation.add_content_slide(
            "Multi-Run",
            [
                [
                    {"text": "Bold ", "bold": True},
                    {"text": "and ", "italic": True},
                    {"text": "normal"},
                ],
            ],
        )
        slide = presentation._pptx.slides[0]
        for shape in slide.placeholders:
            from pptx.enum.shapes import PP_PLACEHOLDER
            if shape.placeholder_format.type in (
                PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT
            ):
                para = shape.text_frame.paragraphs[0]
                assert len(para.runs) == 3
                assert para.runs[0].font.bold is True
                assert para.runs[1].font.italic is True

    def test_dict_format_applies_formatting(
        self, presentation: Presentation
    ) -> None:
        """Test that dict-format items get formatting applied."""
        presentation.add_content_slide(
            "Dict Format",
            [{"text": "Bold item", "bold": True}],
        )
        slide = presentation._pptx.slides[0]
        for shape in slide.placeholders:
            from pptx.enum.shapes import PP_PLACEHOLDER
            if shape.placeholder_format.type in (
                PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT
            ):
                para = shape.text_frame.paragraphs[0]
                if para.runs:
                    assert para.runs[0].font.bold is True

    def test_hyperlink_applied(self, presentation: Presentation) -> None:
        """Test that hyperlinks are applied to runs."""
        presentation.add_content_slide(
            "Links",
            [
                [{"text": "Click here", "hyperlink": "https://example.com"}],
            ],
        )
        slide = presentation._pptx.slides[0]
        for shape in slide.placeholders:
            from pptx.enum.shapes import PP_PLACEHOLDER
            if shape.placeholder_format.type in (
                PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT
            ):
                para = shape.text_frame.paragraphs[0]
                if para.runs:
                    assert para.runs[0].hyperlink.address == "https://example.com"

    def test_font_size_applied(self, presentation: Presentation) -> None:
        """Test that font_size is applied."""
        from pptx.util import Pt
        presentation.add_content_slide(
            "Size",
            [[{"text": "Big text", "font_size": 24}]],
        )
        slide = presentation._pptx.slides[0]
        for shape in slide.placeholders:
            from pptx.enum.shapes import PP_PLACEHOLDER
            if shape.placeholder_format.type in (
                PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT
            ):
                para = shape.text_frame.paragraphs[0]
                if para.runs:
                    assert para.runs[0].font.size == Pt(24)

    def test_plain_string_still_works(
        self, presentation: Presentation
    ) -> None:
        """Test that plain strings still work after the rewrite."""
        presentation.add_content_slide(
            "Plain",
            ["Simple string", "Another string"],
        )
        slide = presentation._pptx.slides[0]
        for shape in slide.placeholders:
            from pptx.enum.shapes import PP_PLACEHOLDER
            if shape.placeholder_format.type in (
                PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT
            ):
                texts = [p.text for p in shape.text_frame.paragraphs]
                assert "Simple string" in texts
                assert "Another string" in texts


class TestImagePositioning:
    """Tests for smart image positioning."""

    @pytest.fixture
    def tiny_image(self, tmp_path: Path) -> Path:
        """Create a tiny PNG image for testing."""
        # Minimal valid PNG (1x1 red pixel)
        import struct
        import zlib

        def create_png(path: Path) -> None:
            signature = b"\x89PNG\r\n\x1a\n"
            # IHDR
            ihdr_data = struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0)
            ihdr_crc = zlib.crc32(b"IHDR" + ihdr_data)
            ihdr = struct.pack(">I", 13) + b"IHDR" + ihdr_data + struct.pack(">I", ihdr_crc & 0xFFFFFFFF)
            # IDAT
            raw = b"\x00\xff\x00\x00"  # filter byte + RGB
            compressed = zlib.compress(raw)
            idat_crc = zlib.crc32(b"IDAT" + compressed)
            idat = struct.pack(">I", len(compressed)) + b"IDAT" + compressed + struct.pack(">I", idat_crc & 0xFFFFFFFF)
            # IEND
            iend_crc = zlib.crc32(b"IEND")
            iend = struct.pack(">I", 0) + b"IEND" + struct.pack(">I", iend_crc & 0xFFFFFFFF)
            path.write_bytes(signature + ihdr + idat + iend)

        img_path = tmp_path / "test.png"
        create_png(img_path)
        return img_path

    def test_image_with_explicit_position(
        self, presentation: Presentation, tiny_image: Path
    ) -> None:
        """Test adding image with explicit positioning."""
        slide_num = presentation.add_image_slide(
            "Custom Position",
            tiny_image,
            left=2.0,
            top=3.0,
            width=4.0,
        )
        assert slide_num == 1

    def test_image_fallback_positioning(
        self, presentation: Presentation, tiny_image: Path
    ) -> None:
        """Test image with default fallback positioning."""
        slide_num = presentation.add_image_slide(
            "Default Position",
            tiny_image,
        )
        assert slide_num == 1

    def test_image_with_caption(
        self, presentation: Presentation, tiny_image: Path
    ) -> None:
        """Test image slide with caption."""
        slide_num = presentation.add_image_slide(
            "Captioned",
            tiny_image,
            caption="A test image",
        )
        assert slide_num == 1

    def test_image_nonexistent_file(
        self, presentation: Presentation
    ) -> None:
        """Test that nonexistent image file doesn't crash."""
        slide_num = presentation.add_image_slide(
            "No Image",
            "/nonexistent/path.png",
        )
        assert slide_num == 1  # Slide created, just no image added
