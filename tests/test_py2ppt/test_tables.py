"""Tests for table slide functionality."""

from pathlib import Path

import pytest
from pptx import Presentation as PptxPresentation

from py2ppt import Template, Presentation, InvalidDataError


@pytest.fixture
def template(tmp_path: Path) -> Template:
    """Create a template for testing."""
    template_path = tmp_path / "template.pptx"
    pres = PptxPresentation()
    pres.slides.add_slide(pres.slide_layouts[0])
    pres.save(str(template_path))
    return Template(template_path)


@pytest.fixture
def presentation(template: Template) -> Presentation:
    """Create a presentation from template."""
    return template.create_presentation()


class TestTableSlide:
    """Tests for add_table_slide."""

    def test_basic_table(self, presentation: Presentation) -> None:
        """Test creating a basic table slide."""
        slide_num = presentation.add_table_slide(
            "Sales Data",
            ["Region", "Q1", "Q2"],
            [
                ["North", "100", "120"],
                ["South", "90", "110"],
            ],
        )
        assert slide_num == 1
        assert presentation.slide_count == 1

    def test_table_data_in_cells(self, presentation: Presentation) -> None:
        """Test that table data is correctly placed in cells."""
        headers = ["Name", "Value"]
        rows = [["Alpha", "10"], ["Beta", "20"]]
        presentation.add_table_slide("Test", headers, rows)

        slide = presentation._pptx.slides[0]
        table_shape = None
        for shape in slide.shapes:
            if shape.has_table:
                table_shape = shape
                break

        assert table_shape is not None
        table = table_shape.table
        assert len(table.rows) == 3  # 1 header + 2 data rows
        assert len(table.columns) == 2
        assert table.cell(0, 0).text == "Name"
        assert table.cell(0, 1).text == "Value"
        assert table.cell(1, 0).text == "Alpha"
        assert table.cell(1, 1).text == "10"
        assert table.cell(2, 0).text == "Beta"
        assert table.cell(2, 1).text == "20"

    def test_table_numeric_values(self, presentation: Presentation) -> None:
        """Test that numeric values are converted to strings."""
        presentation.add_table_slide(
            "Numbers",
            ["Item", "Count"],
            [["Widgets", 42], ["Gadgets", 99]],
        )
        slide = presentation._pptx.slides[0]
        for shape in slide.shapes:
            if shape.has_table:
                assert shape.table.cell(1, 1).text == "42"
                assert shape.table.cell(2, 1).text == "99"

    def test_table_empty_rows(self, presentation: Presentation) -> None:
        """Test table with no data rows."""
        slide_num = presentation.add_table_slide("Empty", ["A", "B"], [])
        assert slide_num == 1

    def test_table_single_column(self, presentation: Presentation) -> None:
        """Test single-column table."""
        slide_num = presentation.add_table_slide(
            "List", ["Items"], [["One"], ["Two"], ["Three"]]
        )
        assert slide_num == 1

    def test_table_custom_col_widths(self, presentation: Presentation) -> None:
        """Test table with custom column widths."""
        slide_num = presentation.add_table_slide(
            "Custom Widths",
            ["Name", "Description"],
            [["A", "Alpha"]],
            col_widths=[2.0, 7.0],
        )
        assert slide_num == 1

    def test_table_theme_style(self, presentation: Presentation) -> None:
        """Test table with theme styling."""
        slide_num = presentation.add_table_slide(
            "Themed",
            ["Col1", "Col2"],
            [["a", "b"]],
            style="theme",
        )
        assert slide_num == 1

    def test_table_plain_style(self, presentation: Presentation) -> None:
        """Test table with plain styling."""
        slide_num = presentation.add_table_slide(
            "Plain",
            ["Col1", "Col2"],
            [["a", "b"]],
            style="plain",
        )
        assert slide_num == 1

    def test_table_striped_style(self, presentation: Presentation) -> None:
        """Test table with striped styling."""
        slide_num = presentation.add_table_slide(
            "Striped",
            ["Col1", "Col2"],
            [["a", "b"], ["c", "d"], ["e", "f"]],
            style="striped",
        )
        assert slide_num == 1

    def test_table_row_mismatch_error(self, presentation: Presentation) -> None:
        """Test that mismatched row/header lengths raise InvalidDataError."""
        with pytest.raises(InvalidDataError) as exc_info:
            presentation.add_table_slide(
                "Bad Data",
                ["A", "B", "C"],
                [["x", "y"]],  # 2 columns, headers has 3
            )
        assert "TABLE_ROW_MISMATCH" in exc_info.value.code
        d = exc_info.value.to_dict()
        assert "suggestion" in d

    def test_table_via_add_slide(self, presentation: Presentation) -> None:
        """Test creating a table via the add_slide auto-router."""
        slide_num = presentation.add_slide(
            content_type="table",
            title="Auto Table",
            headers=["X", "Y"],
            rows=[["1", "2"]],
        )
        assert slide_num == 1

    def test_table_save_and_reload(
        self, presentation: Presentation, tmp_path: Path
    ) -> None:
        """Test that table slides survive save/reload."""
        presentation.add_table_slide(
            "Persist",
            ["A", "B"],
            [["1", "2"]],
        )
        output = tmp_path / "table_output.pptx"
        presentation.save(output)

        loaded = PptxPresentation(str(output))
        assert len(loaded.slides) == 1
        # Verify table exists
        found_table = False
        for shape in loaded.slides[0].shapes:
            if shape.has_table:
                found_table = True
        assert found_table
