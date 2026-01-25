"""Tests for Template class."""

import tempfile
from pathlib import Path

import pytest
from pptx import Presentation as PptxPresentation

from aipptx import Template


@pytest.fixture
def blank_template(tmp_path: Path) -> Path:
    """Create a blank presentation to use as template."""
    template_path = tmp_path / "template.pptx"
    pres = PptxPresentation()
    pres.slides.add_slide(pres.slide_layouts[0])
    pres.save(str(template_path))
    return template_path


class TestTemplate:
    """Tests for the Template class."""

    def test_load_template(self, blank_template: Path) -> None:
        """Test loading a template."""
        template = Template(blank_template)
        assert template.path == blank_template

    def test_template_not_found(self, tmp_path: Path) -> None:
        """Test error when template doesn't exist."""
        with pytest.raises(FileNotFoundError):
            Template(tmp_path / "nonexistent.pptx")

    def test_describe_returns_list(self, blank_template: Path) -> None:
        """Test that describe() returns a list of layout descriptions."""
        template = Template(blank_template)
        layouts = template.describe()

        assert isinstance(layouts, list)
        assert len(layouts) > 0

        # Check structure of first layout
        first = layouts[0]
        assert "name" in first
        assert "index" in first
        assert "type" in first
        assert "description" in first
        assert "placeholders" in first
        assert "best_for" in first

    def test_get_colors(self, blank_template: Path) -> None:
        """Test getting theme colors."""
        template = Template(blank_template)
        colors = template.colors

        assert isinstance(colors, dict)
        # Should have some color entries
        if colors:
            for name, value in colors.items():
                assert isinstance(name, str)
                assert value.startswith("#")

    def test_get_fonts(self, blank_template: Path) -> None:
        """Test getting theme fonts."""
        template = Template(blank_template)
        fonts = template.fonts

        assert isinstance(fonts, dict)
        assert "heading" in fonts
        assert "body" in fonts

    def test_describe_as_text(self, blank_template: Path) -> None:
        """Test getting text description."""
        template = Template(blank_template)
        text = template.describe_as_text()

        assert isinstance(text, str)
        assert "Template:" in text
        assert "Layouts:" in text

    def test_get_layout_by_index(self, blank_template: Path) -> None:
        """Test getting layout by index."""
        template = Template(blank_template)
        layout = template.get_layout(0)

        assert layout is not None
        assert layout.index == 0

    def test_get_layout_by_name(self, blank_template: Path) -> None:
        """Test getting layout by name."""
        template = Template(blank_template)
        layouts = template.describe()

        if layouts:
            name = layouts[0]["name"]
            layout = template.get_layout(name)
            assert layout is not None

    def test_get_layout_names(self, blank_template: Path) -> None:
        """Test getting all layout names."""
        template = Template(blank_template)
        names = template.get_layout_names()

        assert isinstance(names, list)
        assert len(names) > 0
        assert all(isinstance(n, str) for n in names)

    def test_recommend_layout(self, blank_template: Path) -> None:
        """Test layout recommendations."""
        template = Template(blank_template)
        recs = template.recommend_layout("bullets")

        assert isinstance(recs, list)
        # May not have matching layouts in blank template
        for rec in recs:
            assert "name" in rec
            assert "index" in rec
            assert "confidence" in rec
            assert "reason" in rec
            assert 0 <= rec["confidence"] <= 1

    def test_create_presentation(self, blank_template: Path) -> None:
        """Test creating a presentation from template."""
        template = Template(blank_template)
        pres = template.create_presentation()

        from aipptx import Presentation
        assert isinstance(pres, Presentation)
        assert pres.slide_count == 0  # Template slides are removed

    def test_repr(self, blank_template: Path) -> None:
        """Test string representation."""
        template = Template(blank_template)
        repr_str = repr(template)

        assert "Template" in repr_str
        assert "template.pptx" in repr_str


class TestTemplateWithRealFile:
    """Tests with the AWS template if available."""

    @pytest.fixture
    def aws_template(self) -> Path | None:
        """Get the AWS template if it exists."""
        path = Path("/Users/user/Documents/py2ppt/AWStempate.pptx")
        if path.exists():
            return path
        return None

    def test_load_aws_template(self, aws_template: Path | None) -> None:
        """Test loading the AWS template."""
        if aws_template is None:
            pytest.skip("AWS template not available")

        template = Template(aws_template)
        layouts = template.describe()

        # AWS template should have multiple layouts
        assert len(layouts) > 1

        # Print layouts for inspection
        for layout in layouts:
            print(f"\n{layout['name']} ({layout['type']}):")
            print(f"  Description: {layout['description']}")
            print(f"  Placeholders: {list(layout['placeholders'].keys())}")
            print(f"  Best for: {layout['best_for']}")

    def test_aws_template_colors(self, aws_template: Path | None) -> None:
        """Test getting colors from AWS template."""
        if aws_template is None:
            pytest.skip("AWS template not available")

        template = Template(aws_template)
        colors = template.colors

        # AWS template should have theme colors
        assert len(colors) > 0
        print("\nTheme colors:")
        for name, value in colors.items():
            print(f"  {name}: {value}")
