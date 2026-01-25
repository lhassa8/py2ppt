"""Tests for Presentation class."""

from pathlib import Path

import pytest
from pptx import Presentation as PptxPresentation

from aipptx import Template, Presentation


@pytest.fixture
def template(tmp_path: Path) -> Template:
    """Create a template for testing."""
    template_path = tmp_path / "template.pptx"
    pres = PptxPresentation()
    pres.slides.add_slide(pres.slide_layouts[0])
    pres.save(str(template_path))
    return Template(template_path)


@pytest.fixture
def presentation(template: Template) -> Presentation:
    """Create a presentation from template."""
    return template.create_presentation()


class TestPresentation:
    """Tests for the Presentation class."""

    def test_create_presentation(self, template: Template) -> None:
        """Test creating a presentation."""
        pres = template.create_presentation()
        assert isinstance(pres, Presentation)
        assert pres.slide_count == 0

    def test_add_title_slide(self, presentation: Presentation) -> None:
        """Test adding a title slide."""
        slide_num = presentation.add_title_slide("Test Title", "Subtitle")

        assert slide_num == 1
        assert presentation.slide_count == 1

    def test_add_content_slide(self, presentation: Presentation) -> None:
        """Test adding a content slide."""
        slide_num = presentation.add_content_slide(
            "Content Slide",
            ["Point 1", "Point 2", "Point 3"]
        )

        assert slide_num == 1
        assert presentation.slide_count == 1

    def test_add_content_slide_with_levels(self, presentation: Presentation) -> None:
        """Test adding a content slide with nested bullets."""
        slide_num = presentation.add_content_slide(
            "Nested Content",
            ["Main point", "Sub point 1", "Sub point 2", "Another main"],
            levels=[0, 1, 1, 0]
        )

        assert slide_num == 1
        assert presentation.slide_count == 1

    def test_add_section_slide(self, presentation: Presentation) -> None:
        """Test adding a section slide."""
        slide_num = presentation.add_section_slide("Section Title")

        assert slide_num == 1
        assert presentation.slide_count == 1

    def test_add_two_column_slide(self, presentation: Presentation) -> None:
        """Test adding a two-column slide."""
        slide_num = presentation.add_two_column_slide(
            "Two Columns",
            ["Left 1", "Left 2"],
            ["Right 1", "Right 2"]
        )

        assert slide_num == 1
        assert presentation.slide_count == 1

    def test_add_comparison_slide(self, presentation: Presentation) -> None:
        """Test adding a comparison slide."""
        slide_num = presentation.add_comparison_slide(
            "Comparison",
            "Before", ["Old way 1", "Old way 2"],
            "After", ["New way 1", "New way 2"]
        )

        assert slide_num == 1
        assert presentation.slide_count == 1

    def test_add_blank_slide(self, presentation: Presentation) -> None:
        """Test adding a blank slide."""
        slide_num = presentation.add_blank_slide()

        assert slide_num == 1
        assert presentation.slide_count == 1

    def test_add_slide_auto(self, presentation: Presentation) -> None:
        """Test add_slide with auto layout."""
        slide_num = presentation.add_slide(
            content_type="content",
            title="Auto Slide",
            content=["Point 1", "Point 2"]
        )

        assert slide_num == 1

    def test_add_multiple_slides(self, presentation: Presentation) -> None:
        """Test adding multiple slides."""
        presentation.add_title_slide("Title", "Subtitle")
        presentation.add_content_slide("Content", ["Point"])
        presentation.add_section_slide("Section")

        assert presentation.slide_count == 3

    def test_set_notes(self, presentation: Presentation) -> None:
        """Test setting speaker notes."""
        presentation.add_title_slide("Title")
        presentation.set_notes(1, "These are speaker notes.")

        # Verify notes were set
        slide = presentation._pptx.slides[0]
        notes = slide.notes_slide.notes_text_frame.text
        assert "speaker notes" in notes.lower()

    def test_save(self, presentation: Presentation, tmp_path: Path) -> None:
        """Test saving a presentation."""
        presentation.add_title_slide("Test")

        output_path = tmp_path / "output.pptx"
        presentation.save(output_path)

        assert output_path.exists()

        # Verify saved file can be opened
        loaded = PptxPresentation(str(output_path))
        assert len(loaded.slides) == 1

    def test_repr(self, presentation: Presentation) -> None:
        """Test string representation."""
        presentation.add_title_slide("Test")
        repr_str = repr(presentation)

        assert "Presentation" in repr_str
        assert "1 slides" in repr_str

    def test_template_property(self, presentation: Presentation, template: Template) -> None:
        """Test accessing the template property."""
        assert presentation.template == template


class TestContentFormatting:
    """Tests for content formatting in presentations."""

    def test_rich_text_content(self, presentation: Presentation) -> None:
        """Test adding content with rich text."""
        slide_num = presentation.add_content_slide(
            "Rich Text",
            [
                [{"text": "Bold: ", "bold": True}, {"text": "normal"}],
                "Plain bullet",
            ]
        )

        assert slide_num == 1

    def test_tuple_format_content(self, presentation: Presentation) -> None:
        """Test adding content with tuple format."""
        slide_num = presentation.add_content_slide(
            "Tuple Format",
            [
                "Main point",
                ("Sub point", 1),
                ("Deep sub", 2),
            ]
        )

        assert slide_num == 1

    def test_dict_format_content(self, presentation: Presentation) -> None:
        """Test adding content with dict format."""
        slide_num = presentation.add_content_slide(
            "Dict Format",
            [
                {"text": "Formatted", "bold": True},
                "Plain",
            ]
        )

        assert slide_num == 1


class TestWithRealTemplate:
    """Tests with the AWS template if available."""

    @pytest.fixture
    def aws_template(self) -> Template | None:
        """Get the AWS template if it exists."""
        path = Path("/Users/user/Documents/py2ppt/AWStempate.pptx")
        if path.exists():
            return Template(path)
        return None

    def test_create_full_presentation(
        self, aws_template: Template | None, tmp_path: Path
    ) -> None:
        """Test creating a full presentation with AWS template."""
        if aws_template is None:
            pytest.skip("AWS template not available")

        pres = aws_template.create_presentation()

        # Add various slide types
        pres.add_title_slide("Q4 Business Review", "January 2025")

        pres.add_content_slide("Key Highlights", [
            "Revenue exceeded targets by 15%",
            "Customer satisfaction at all-time high",
            "Three new product launches completed"
        ])

        pres.add_comparison_slide(
            "Before vs After",
            "Legacy System", ["Slow", "Manual", "Error-prone"],
            "New Platform", ["Fast", "Automated", "Reliable"]
        )

        pres.add_section_slide("Next Steps")

        pres.add_content_slide("Action Items", [
            "Complete migration by Q1",
            "Training for all teams",
            "Monitor performance metrics"
        ])

        # Save and verify
        output_path = tmp_path / "test_presentation.pptx"
        pres.save(output_path)

        assert output_path.exists()
        assert pres.slide_count == 5

        # Verify file can be opened
        loaded = PptxPresentation(str(output_path))
        assert len(loaded.slides) == 5
