"""Presentation class with semantic slide management.

Provides high-level methods for creating presentations with
AI-friendly, intent-based APIs.
"""

from __future__ import annotations

from pathlib import Path
from typing import TYPE_CHECKING, Any

from pptx import Presentation as PptxPresentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import PP_PLACEHOLDER, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
from pptx.oxml.ns import qn

from .errors import (
    SlideNotFoundError,
    LayoutNotFoundError,
    InvalidDataError,
)
from .formatting import format_for_py2ppt, parse_content
from .layout import LayoutType

if TYPE_CHECKING:
    from .template import Template

# Chart type mapping
_CHART_TYPE_MAP = {
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "line": XL_CHART_TYPE.LINE,
    "pie": XL_CHART_TYPE.PIE,
    "doughnut": XL_CHART_TYPE.DOUGHNUT,
}


class Presentation:
    """AI-friendly presentation with semantic slide methods.

    This class wraps python-pptx's Presentation to provide high-level,
    intent-based methods for creating slides. Instead of dealing
    with placeholder indices and layout details, you can simply
    call methods like add_title_slide() or add_comparison_slide().

    Example:
        >>> template = Template("corporate.pptx")
        >>> pres = template.create_presentation()
        >>>
        >>> pres.add_title_slide("Q4 Review", "January 2025")
        >>> pres.add_content_slide("Key Points", [
        ...     "Revenue up 20%",
        ...     "New markets opened",
        ...     "Customer satisfaction high"
        ... ])
        >>> pres.save("output.pptx")
    """

    def __init__(self, template: "Template") -> None:
        """Initialize from a Template.

        Use Template.create_presentation() instead of this constructor.
        """
        self._template = template
        self._layouts = template._layouts

        # Create python-pptx presentation from template
        self._pptx = PptxPresentation(template.path)

        # Remove existing slides (keep only layouts/masters)
        while len(self._pptx.slides) > 0:
            rId = self._pptx.slides._sldIdLst[0].rId
            self._pptx.part.drop_rel(rId)
            del self._pptx.slides._sldIdLst[0]

    @property
    def slide_count(self) -> int:
        """Get the current number of slides."""
        return len(self._pptx.slides)

    @property
    def template(self) -> "Template":
        """Get the template this presentation was created from."""
        return self._template

    # --- Private helpers ---

    def _validate_slide_number(self, n: int) -> None:
        """Validate a slide number, raising SlideNotFoundError if invalid."""
        count = len(self._pptx.slides)
        if not isinstance(n, int) or n < 1 or n > count:
            raise SlideNotFoundError(
                f"Slide {n} does not exist. Presentation has {count} slide(s).",
                suggestion=(
                    f"Use a slide number between 1 and {count}."
                    if count > 0
                    else "Add slides first."
                ),
                code="SLIDE_NOT_FOUND",
            )

    def _find_layout_by_type(self, layout_type: LayoutType) -> int:
        """Find the best layout index for a layout type."""
        for layout in self._layouts:
            if layout.layout_type == layout_type:
                return layout.index
        # Fall back to first layout
        return 0

    def _find_layout(self, layout: str | int | None, layout_type: LayoutType) -> int:
        """Find layout by name, index, or type.

        Raises LayoutNotFoundError for explicit names/indices that don't match.
        Falls back silently for None/"auto".
        """
        if layout is None or layout == "auto":
            return self._find_layout_by_type(layout_type)

        if isinstance(layout, int):
            if layout < 0 or layout >= len(self._pptx.slide_layouts):
                raise LayoutNotFoundError(
                    f"Layout index {layout} is out of range.",
                    suggestion=f"Available indices: 0-{len(self._pptx.slide_layouts) - 1}.",
                    code="LAYOUT_NOT_FOUND",
                )
            return layout

        # Fuzzy name match
        layout_lower = layout.lower()
        for l in self._layouts:
            if layout_lower in l.name.lower() or l.name.lower() in layout_lower:
                return l.index

        raise LayoutNotFoundError(
            f"No layout matching '{layout}' found.",
            suggestion=f"Available layouts: {', '.join(l.name for l in self._layouts)}.",
            code="LAYOUT_NOT_FOUND",
        )

    def _get_placeholder(self, slide, ph_type):
        """Get a placeholder by type from a slide."""
        for shape in slide.placeholders:
            if shape.placeholder_format.type == ph_type:
                return shape
        return None

    def _set_text_frame(self, shape, text: str) -> None:
        """Set text in a shape's text frame."""
        if shape is None:
            return
        tf = shape.text_frame
        tf.clear()
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        p = tf.paragraphs[0]
        p.text = text

    def _apply_run_formatting(self, run, fmt: dict) -> None:
        """Apply formatting from a dict to a python-pptx Run.

        Supports: bold, italic, underline, font_size, font_family, color, hyperlink.
        """
        if fmt.get("bold"):
            run.font.bold = True
        if fmt.get("italic"):
            run.font.italic = True
        if fmt.get("underline"):
            run.font.underline = True
        if fmt.get("font_size"):
            run.font.size = Pt(fmt["font_size"])
        if fmt.get("font_family"):
            run.font.name = fmt["font_family"]
        if fmt.get("color"):
            color_hex = fmt["color"].lstrip("#")
            if len(color_hex) == 6:
                run.font.color.rgb = RGBColor(
                    int(color_hex[:2], 16),
                    int(color_hex[2:4], 16),
                    int(color_hex[4:6], 16),
                )
        if fmt.get("hyperlink"):
            run.hyperlink.address = fmt["hyperlink"]

    def _set_paragraph_runs(self, p, segments: list) -> None:
        """Set formatted runs on a paragraph, replacing any existing runs."""
        # Remove existing <a:r> elements from the paragraph XML
        for r_elem in list(p._p):
            if r_elem.tag == qn("a:r"):
                p._p.remove(r_elem)

        for seg in segments:
            run = p.add_run()
            if isinstance(seg, dict):
                run.text = seg.get("text", "")
                self._apply_run_formatting(run, seg)
            elif isinstance(seg, str):
                run.text = seg

    def _set_body_content(
        self, shape, content: list, levels: list[int] | None = None
    ) -> None:
        """Set bullet content in a body placeholder with full rich text support."""
        if shape is None:
            return

        tf = shape.text_frame
        tf.clear()
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        for i, item in enumerate(content):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            # Set level
            if levels and i < len(levels):
                p.level = levels[i]

            # Set text with formatting
            if isinstance(item, str):
                p.text = item
            elif isinstance(item, list):
                # Rich text: list of dicts/strings with formatting
                self._set_paragraph_runs(p, item)
            elif isinstance(item, dict):
                # Single formatted run
                self._set_paragraph_runs(p, [item])

    def _get_theme_color(self, name: str = "accent1") -> RGBColor:
        """Get a theme color as RGBColor, with fallback."""
        hex_color = self._template.colors.get(name, "#4472C4").lstrip("#")
        return RGBColor(
            int(hex_color[:2], 16),
            int(hex_color[2:4], 16),
            int(hex_color[4:6], 16),
        )

    # --- Slide creation methods ---

    def add_title_slide(
        self,
        title: str,
        subtitle: str = "",
        *,
        layout: str | int | None = None,
    ) -> int:
        """Add a title/cover slide.

        Args:
            title: Main presentation title
            subtitle: Optional subtitle (presenter name, date, etc.)
            layout: Layout name, index, or None for auto-selection

        Returns:
            Slide number of the new slide

        Example:
            >>> pres.add_title_slide("Q4 Business Review", "January 2025")
        """
        layout_idx = self._find_layout(layout, LayoutType.TITLE)
        slide_layout = self._pptx.slide_layouts[layout_idx]
        slide = self._pptx.slides.add_slide(slide_layout)

        # Set title
        title_ph = self._get_placeholder(slide, PP_PLACEHOLDER.TITLE)
        if title_ph is None:
            title_ph = self._get_placeholder(slide, PP_PLACEHOLDER.CENTER_TITLE)
        self._set_text_frame(title_ph, title)

        # Set subtitle
        if subtitle:
            subtitle_ph = self._get_placeholder(slide, PP_PLACEHOLDER.SUBTITLE)
            self._set_text_frame(subtitle_ph, subtitle)

        return len(self._pptx.slides)

    def add_section_slide(
        self,
        title: str,
        subtitle: str = "",
        *,
        layout: str | int | None = None,
    ) -> int:
        """Add a section header/divider slide.

        Args:
            title: Section title
            subtitle: Optional subtitle or description
            layout: Layout name, index, or None for auto-selection

        Returns:
            Slide number of the new slide

        Example:
            >>> pres.add_section_slide("Part 2: Analysis")
        """
        layout_idx = self._find_layout(layout, LayoutType.SECTION)
        slide_layout = self._pptx.slide_layouts[layout_idx]
        slide = self._pptx.slides.add_slide(slide_layout)

        # Set title
        title_ph = self._get_placeholder(slide, PP_PLACEHOLDER.TITLE)
        self._set_text_frame(title_ph, title)

        # Set subtitle if provided
        if subtitle:
            subtitle_ph = self._get_placeholder(slide, PP_PLACEHOLDER.SUBTITLE)
            if subtitle_ph:
                self._set_text_frame(subtitle_ph, subtitle)

        return len(self._pptx.slides)

    def add_content_slide(
        self,
        title: str,
        content: str | list[str | tuple | dict | list],
        *,
        levels: list[int] | None = None,
        layout: str | int | None = None,
        warn_overflow: bool = False,
    ) -> int | dict:
        """Add a content slide with bullets.

        Args:
            title: Slide title
            content: Content as string, list of strings, or rich text
            levels: Optional indent levels for each bullet (0=top level)
            layout: Layout name, index, or None for auto-selection
            warn_overflow: If True, returns a dict with overflow info instead of bare int

        Returns:
            Slide number (int) or overflow info dict if warn_overflow=True

        Example:
            >>> pres.add_content_slide("Key Points", [
            ...     "First point",
            ...     "Second point",
            ...     "Third point"
            ... ])

            >>> # With overflow detection
            >>> result = pres.add_content_slide("Many Points", items, warn_overflow=True)
            >>> if result["overflow"]:
            ...     print(f"Content overflows: {result['item_count']} items, capacity ~{result['estimated_capacity']}")
        """
        layout_idx = self._find_layout(layout, LayoutType.CONTENT)
        slide_layout = self._pptx.slide_layouts[layout_idx]
        slide = self._pptx.slides.add_slide(slide_layout)

        # Set title
        title_ph = self._get_placeholder(slide, PP_PLACEHOLDER.TITLE)
        self._set_text_frame(title_ph, title)

        # Parse and format content
        paragraphs = parse_content(content, levels)
        formatted_content, formatted_levels = format_for_py2ppt(paragraphs)

        # Set body
        body_ph = self._get_placeholder(slide, PP_PLACEHOLDER.BODY)
        if body_ph is None:
            body_ph = self._get_placeholder(slide, PP_PLACEHOLDER.OBJECT)
        self._set_body_content(body_ph, formatted_content, formatted_levels)

        slide_num = len(self._pptx.slides)

        if warn_overflow:
            # Estimate capacity from placeholder height (~0.4 inches per bullet)
            capacity = 6  # default
            if body_ph is not None:
                height_inches = body_ph.height / 914400  # EMU to inches
                capacity = max(1, int(height_inches / 0.4))
            item_count = len(paragraphs)
            return {
                "slide_number": slide_num,
                "overflow": item_count > capacity,
                "item_count": item_count,
                "estimated_capacity": capacity,
            }

        return slide_num

    def add_two_column_slide(
        self,
        title: str,
        left_content: list[str | tuple | dict | list],
        right_content: list[str | tuple | dict | list],
        *,
        left_levels: list[int] | None = None,
        right_levels: list[int] | None = None,
        layout: str | int | None = None,
    ) -> int:
        """Add a two-column content slide.

        Args:
            title: Slide title
            left_content: Content for left column
            right_content: Content for right column
            left_levels: Indent levels for left content
            right_levels: Indent levels for right content
            layout: Layout name, index, or None for auto-selection

        Returns:
            Slide number of the new slide

        Example:
            >>> pres.add_two_column_slide(
            ...     "Topics",
            ...     ["Point A", "Point B"],
            ...     ["Point X", "Point Y"]
            ... )
        """
        layout_idx = self._find_layout(layout, LayoutType.TWO_COLUMN)
        slide_layout = self._pptx.slide_layouts[layout_idx]
        slide = self._pptx.slides.add_slide(slide_layout)

        # Set title
        title_ph = self._get_placeholder(slide, PP_PLACEHOLDER.TITLE)
        self._set_text_frame(title_ph, title)

        # Find body placeholders (there should be two)
        body_phs = [
            shape
            for shape in slide.placeholders
            if shape.placeholder_format.type
            in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT)
        ]
        body_phs.sort(key=lambda s: s.left)  # Sort by x position

        # Format content
        left_paragraphs = parse_content(left_content, left_levels)
        left_formatted, left_lvls = format_for_py2ppt(left_paragraphs)

        right_paragraphs = parse_content(right_content, right_levels)
        right_formatted, right_lvls = format_for_py2ppt(right_paragraphs)

        # Set left content
        if len(body_phs) >= 1:
            self._set_body_content(body_phs[0], left_formatted, left_lvls)

        # Set right content
        if len(body_phs) >= 2:
            self._set_body_content(body_phs[1], right_formatted, right_lvls)

        return len(self._pptx.slides)

    def add_comparison_slide(
        self,
        title: str,
        left_heading: str,
        left_content: list[str | tuple | dict | list],
        right_heading: str,
        right_content: list[str | tuple | dict | list],
        *,
        layout: str | int | None = None,
    ) -> int:
        """Add a comparison slide with two titled columns.

        Ideal for before/after, pros/cons, or versus comparisons.

        Args:
            title: Slide title
            left_heading: Heading for left column
            left_content: Bullet points for left column
            right_heading: Heading for right column
            right_content: Bullet points for right column
            layout: Layout name, index, or None for auto-selection

        Returns:
            Slide number of the new slide

        Example:
            >>> pres.add_comparison_slide(
            ...     "Before vs After",
            ...     "Legacy System",
            ...     ["Slow", "Manual", "Error-prone"],
            ...     "New Platform",
            ...     ["Fast", "Automated", "Reliable"]
            ... )
        """
        # Try to find a comparison layout, fall back to two-column
        layout_idx = self._find_layout(layout, LayoutType.COMPARISON)
        slide_layout = self._pptx.slide_layouts[layout_idx]
        slide = self._pptx.slides.add_slide(slide_layout)

        # Set title
        title_ph = self._get_placeholder(slide, PP_PLACEHOLDER.TITLE)
        self._set_text_frame(title_ph, title)

        # Find body placeholders
        body_phs = [
            shape
            for shape in slide.placeholders
            if shape.placeholder_format.type
            in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT)
        ]
        # Sort by position: top row first, then left to right
        body_phs.sort(key=lambda s: (s.top, s.left))

        if len(body_phs) >= 4:
            # True comparison layout: heading, content, heading, content
            self._set_text_frame(body_phs[0], left_heading)
            left_paragraphs = parse_content(left_content)
            left_formatted, left_lvls = format_for_py2ppt(left_paragraphs)
            self._set_body_content(body_phs[2], left_formatted, left_lvls)

            self._set_text_frame(body_phs[1], right_heading)
            right_paragraphs = parse_content(right_content)
            right_formatted, right_lvls = format_for_py2ppt(right_paragraphs)
            self._set_body_content(body_phs[3], right_formatted, right_lvls)

        elif len(body_phs) >= 2:
            # Two-column layout - combine heading with content
            left_combined = [left_heading] + list(left_content)
            right_combined = [right_heading] + list(right_content)

            left_levels = [0] + [1] * len(left_content)
            right_levels = [0] + [1] * len(right_content)

            left_paragraphs = parse_content(left_combined, left_levels)
            right_paragraphs = parse_content(right_combined, right_levels)

            left_formatted, left_lvls = format_for_py2ppt(left_paragraphs)
            right_formatted, right_lvls = format_for_py2ppt(right_paragraphs)

            # Sort by x position for left/right
            body_phs.sort(key=lambda s: s.left)
            self._set_body_content(body_phs[0], left_formatted, left_lvls)
            self._set_body_content(body_phs[1], right_formatted, right_lvls)

        elif len(body_phs) >= 1:
            # Single body - combine all
            combined = [
                left_heading,
                *list(left_content),
                "",
                right_heading,
                *list(right_content),
            ]
            combined_levels = (
                [0]
                + [1] * len(left_content)
                + [0]
                + [0]
                + [1] * len(right_content)
            )
            paragraphs = parse_content(combined, combined_levels)
            formatted, lvls = format_for_py2ppt(paragraphs)
            self._set_body_content(body_phs[0], formatted, lvls)

        return len(self._pptx.slides)

    def add_image_slide(
        self,
        title: str,
        image_path: str | Path,
        caption: str = "",
        *,
        layout: str | int | None = None,
        left: float | None = None,
        top: float | None = None,
        width: float | None = None,
        height: float | None = None,
    ) -> int:
        """Add a slide with an image.

        Image placement priority:
        1. PICTURE placeholder (uses insert_picture for best fit)
        2. Explicit left/top/width/height kwargs (in inches)
        3. BODY/OBJECT placeholder bounds as fallback
        4. Default position (1in, 2in, width=5in)

        Args:
            title: Slide title
            image_path: Path to the image file
            caption: Optional caption text
            layout: Layout name, index, or None for auto-selection
            left: Image left position in inches
            top: Image top position in inches
            width: Image width in inches
            height: Image height in inches

        Returns:
            Slide number of the new slide

        Example:
            >>> pres.add_image_slide("Product Photo", "product.png", "Our flagship product")
            >>> pres.add_image_slide("Custom Position", "img.png", left=2, top=3, width=6)
        """
        layout_idx = self._find_layout(layout, LayoutType.IMAGE_CONTENT)
        slide_layout = self._pptx.slide_layouts[layout_idx]
        slide = self._pptx.slides.add_slide(slide_layout)

        # Set title
        title_ph = self._get_placeholder(slide, PP_PLACEHOLDER.TITLE)
        self._set_text_frame(title_ph, title)

        # Add image
        image_path = Path(image_path)
        if image_path.exists():
            added = False

            # 1. Try PICTURE placeholder
            pic_ph = self._get_placeholder(slide, PP_PLACEHOLDER.PICTURE)
            if pic_ph is not None:
                pic_ph.insert_picture(str(image_path))
                added = True

            # 2. Explicit dimensions
            if not added and any(
                v is not None for v in [left, top, width, height]
            ):
                pic_left = Inches(left) if left is not None else Inches(1)
                pic_top = Inches(top) if top is not None else Inches(2)
                pic_width = Inches(width) if width is not None else None
                pic_height = Inches(height) if height is not None else None
                slide.shapes.add_picture(
                    str(image_path),
                    pic_left,
                    pic_top,
                    width=pic_width,
                    height=pic_height,
                )
                added = True

            # 3. BODY/OBJECT placeholder bounds
            if not added:
                body_ph = self._get_placeholder(slide, PP_PLACEHOLDER.BODY)
                if body_ph is None:
                    body_ph = self._get_placeholder(
                        slide, PP_PLACEHOLDER.OBJECT
                    )
                if body_ph is not None:
                    slide.shapes.add_picture(
                        str(image_path),
                        body_ph.left,
                        body_ph.top,
                        width=body_ph.width,
                    )
                else:
                    # 4. Default fallback
                    slide.shapes.add_picture(
                        str(image_path),
                        Inches(1),
                        Inches(2),
                        width=Inches(5),
                    )

        # Set caption if provided
        if caption:
            body_ph = self._get_placeholder(slide, PP_PLACEHOLDER.BODY)
            if body_ph:
                self._set_text_frame(body_ph, caption)

        return len(self._pptx.slides)

    def add_blank_slide(self, layout: str | int | None = None) -> int:
        """Add a blank slide.

        Args:
            layout: Layout name, index, or None for auto-selection

        Returns:
            Slide number of the new slide

        Example:
            >>> slide_num = pres.add_blank_slide()
        """
        layout_idx = self._find_layout(layout, LayoutType.BLANK)
        slide_layout = self._pptx.slide_layouts[layout_idx]
        self._pptx.slides.add_slide(slide_layout)
        return len(self._pptx.slides)

    def add_table_slide(
        self,
        title: str,
        headers: list[str],
        rows: list[list[Any]],
        *,
        col_widths: list[float] | None = None,
        style: str = "theme",
        layout: str | int | None = None,
    ) -> int:
        """Add a slide with a table.

        Args:
            title: Slide title
            headers: Column header labels
            rows: List of rows, each a list of cell values
            col_widths: Optional column widths in inches (auto-calculated if omitted)
            style: Table style - "theme" (default), "plain", or "striped"
            layout: Layout name, index, or None for auto-selection

        Returns:
            Slide number of the new slide

        Example:
            >>> pres.add_table_slide(
            ...     "Sales Data",
            ...     ["Region", "Q1", "Q2", "Q3"],
            ...     [["North", 100, 120, 130],
            ...      ["South", 90, 110, 125]],
            ... )
        """
        # Validate row/header match
        for i, row in enumerate(rows):
            if len(row) != len(headers):
                raise InvalidDataError(
                    f"Row {i} has {len(row)} columns but headers has {len(headers)}.",
                    suggestion="Ensure all rows have the same number of columns as headers.",
                    code="TABLE_ROW_MISMATCH",
                )

        layout_idx = self._find_layout(layout, LayoutType.TITLE_ONLY)
        slide_layout = self._pptx.slide_layouts[layout_idx]
        slide = self._pptx.slides.add_slide(slide_layout)

        # Title
        title_ph = self._get_placeholder(slide, PP_PLACEHOLDER.TITLE)
        self._set_text_frame(title_ph, title)

        # Table dimensions
        num_rows = len(rows) + 1  # +1 for header
        num_cols = len(headers)
        tbl_left = Inches(0.5)
        tbl_top = Inches(1.8)
        tbl_width = Inches(9.0)
        tbl_height = Inches(min(0.4 * num_rows, 5.2))

        table_shape = slide.shapes.add_table(
            num_rows, num_cols, tbl_left, tbl_top, tbl_width, tbl_height
        )
        table = table_shape.table

        # Set column widths
        if col_widths:
            for c, w in enumerate(col_widths):
                if c < num_cols:
                    table.columns[c].width = Inches(w)
        else:
            # Auto-calculate proportional to content length
            max_lengths = []
            for c in range(num_cols):
                max_len = len(str(headers[c]))
                for row in rows:
                    max_len = max(max_len, len(str(row[c])))
                max_lengths.append(max(max_len, 1))
            total_len = sum(max_lengths)
            for c in range(num_cols):
                table.columns[c].width = int(
                    tbl_width * max_lengths[c] / total_len
                )

        # Set headers
        for c, header in enumerate(headers):
            table.cell(0, c).text = str(header)

        # Set data rows
        for r, row in enumerate(rows):
            for c, value in enumerate(row):
                table.cell(r + 1, c).text = str(value)

        # Apply styling
        accent_color = self._get_theme_color("accent1")

        if style in ("theme", "striped"):
            # Style header row
            for c in range(num_cols):
                cell = table.cell(0, c)
                cell.fill.solid()
                cell.fill.fore_color.rgb = accent_color
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.bold = True
                        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

            # Striped: alternate row backgrounds
            if style == "striped":
                for r in range(len(rows)):
                    if r % 2 == 0:
                        for c in range(num_cols):
                            cell = table.cell(r + 1, c)
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = RGBColor(
                                0xF2, 0xF2, 0xF2
                            )

        # "plain" style: no special formatting

        return len(self._pptx.slides)

    def add_chart_slide(
        self,
        title: str,
        chart_type: str,
        data: dict[str, Any],
        *,
        layout: str | int | None = None,
    ) -> int:
        """Add a slide with a chart.

        Args:
            title: Slide title
            chart_type: One of "bar", "column", "line", "pie", "doughnut"
            data: Chart data dict. Format depends on chart type:
                  Bar/column/line: {"categories": [...], "series": [{"name": "...", "values": [...]}]}
                  Pie/doughnut: {"categories": [...], "values": [...]}
            layout: Layout name, index, or None for auto-selection

        Returns:
            Slide number of the new slide

        Example:
            >>> pres.add_chart_slide("Revenue", "column", {
            ...     "categories": ["Q1", "Q2", "Q3"],
            ...     "series": [{"name": "2024", "values": [10, 20, 30]}]
            ... })
        """
        # Validate chart type
        chart_type_lower = chart_type.lower()
        if chart_type_lower not in _CHART_TYPE_MAP:
            raise InvalidDataError(
                f"Unknown chart type '{chart_type}'.",
                suggestion=f"Supported types: {', '.join(_CHART_TYPE_MAP.keys())}.",
                code="INVALID_CHART_TYPE",
            )

        # Validate data
        if "categories" not in data:
            raise InvalidDataError(
                "Chart data must include 'categories' key.",
                suggestion=(
                    "Provide {'categories': [...], 'series': [...]} "
                    "or {'categories': [...], 'values': [...]}."
                ),
                code="MISSING_CHART_DATA",
            )

        xl_chart_type = _CHART_TYPE_MAP[chart_type_lower]

        # Build chart data
        chart_data = CategoryChartData()
        chart_data.categories = data["categories"]

        if chart_type_lower in ("pie", "doughnut"):
            if "values" not in data:
                raise InvalidDataError(
                    f"'{chart_type}' chart requires 'values' key in data.",
                    suggestion="Provide {'categories': [...], 'values': [...]}.",
                    code="MISSING_CHART_VALUES",
                )
            chart_data.add_series("Values", data["values"])
        else:
            if "series" not in data:
                raise InvalidDataError(
                    f"'{chart_type}' chart requires 'series' key in data.",
                    suggestion=(
                        "Provide {'categories': [...], "
                        "'series': [{'name': '...', 'values': [...]}]}."
                    ),
                    code="MISSING_CHART_SERIES",
                )
            for series in data["series"]:
                chart_data.add_series(
                    series.get("name", "Series"),
                    series.get("values", []),
                )

        # Create slide
        layout_idx = self._find_layout(layout, LayoutType.TITLE_ONLY)
        slide_layout = self._pptx.slide_layouts[layout_idx]
        slide = self._pptx.slides.add_slide(slide_layout)

        # Title
        title_ph = self._get_placeholder(slide, PP_PLACEHOLDER.TITLE)
        self._set_text_frame(title_ph, title)

        # Add chart
        chart_left = Inches(1.0)
        chart_top = Inches(1.8)
        chart_width = Inches(8.0)
        chart_height = Inches(5.0)

        chart_shape = slide.shapes.add_chart(
            xl_chart_type,
            chart_left,
            chart_top,
            chart_width,
            chart_height,
            chart_data,
        )
        chart = chart_shape.chart

        # Add legend for multi-series charts
        if chart_type_lower not in ("pie", "doughnut"):
            if len(data.get("series", [])) > 1:
                chart.has_legend = True

        return len(self._pptx.slides)

    def add_slide(
        self,
        layout: str | int = "auto",
        content_type: str = "content",
        title: str = "",
        content: list[str] | None = None,
        **kwargs: Any,
    ) -> int:
        """Add a slide with auto-layout selection.

        A flexible method that chooses the right slide type
        based on content_type and delegates to the appropriate
        specialized method.

        Args:
            layout: Layout name, index, or "auto"
            content_type: Type of content ("title", "content", "comparison",
                          "table", "chart", etc.)
            title: Slide title
            content: Content items
            **kwargs: Additional arguments for the specific slide type

        Returns:
            Slide number of the new slide

        Example:
            >>> pres.add_slide(content_type="table", title="Data",
            ...     headers=["A", "B"], rows=[[1, 2]])
        """
        content_lower = content_type.lower()

        if content_lower in ("title", "cover", "opening"):
            subtitle = kwargs.get("subtitle", "")
            return self.add_title_slide(title, subtitle, layout=layout)

        elif content_lower in ("section", "divider"):
            subtitle = kwargs.get("subtitle", "")
            return self.add_section_slide(title, subtitle, layout=layout)

        elif content_lower in ("comparison", "versus", "vs"):
            return self.add_comparison_slide(
                title,
                kwargs.get("left_heading", "Option A"),
                kwargs.get("left_content", content or []),
                kwargs.get("right_heading", "Option B"),
                kwargs.get("right_content", []),
                layout=layout,
            )

        elif content_lower in ("two_column", "split"):
            return self.add_two_column_slide(
                title,
                kwargs.get("left_content", content or []),
                kwargs.get("right_content", []),
                layout=layout,
            )

        elif content_lower in ("image", "picture"):
            return self.add_image_slide(
                title,
                kwargs.get("image_path", ""),
                kwargs.get("caption", ""),
                layout=layout,
            )

        elif content_lower == "blank":
            return self.add_blank_slide(layout=layout)

        elif content_lower == "table":
            return self.add_table_slide(
                title,
                kwargs.get("headers", []),
                kwargs.get("rows", []),
                col_widths=kwargs.get("col_widths"),
                style=kwargs.get("style", "theme"),
                layout=layout,
            )

        elif content_lower == "chart":
            return self.add_chart_slide(
                title,
                kwargs.get("chart_type", "column"),
                kwargs.get("data", {}),
                layout=layout,
            )

        else:
            # Default to content slide
            return self.add_content_slide(
                title,
                content or [],
                levels=kwargs.get("levels"),
                layout=layout,
            )

    def set_notes(self, slide_number: int, notes: str) -> None:
        """Set speaker notes for a slide.

        Args:
            slide_number: The slide number (1-indexed)
            notes: Notes text (can include newlines)

        Raises:
            SlideNotFoundError: If slide_number is out of range

        Example:
            >>> pres.set_notes(1, "Remember to mention...")
        """
        self._validate_slide_number(slide_number)

        slide = self._pptx.slides[slide_number - 1]
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

    # --- Inspection methods ---

    def describe_slide(self, n: int) -> dict[str, Any]:
        """Describe a single slide's contents.

        Args:
            n: Slide number (1-indexed)

        Returns:
            Dict with: slide_number, layout, title, content, shapes, notes,
            has_title, has_content, has_table, has_chart, has_image

        Raises:
            SlideNotFoundError: If slide number is out of range

        Example:
            >>> info = pres.describe_slide(1)
            >>> print(info["title"], info["has_table"])
        """
        self._validate_slide_number(n)
        slide = self._pptx.slides[n - 1]

        layout_name = slide.slide_layout.name

        # Find title
        title = ""
        has_title = False
        for shape in slide.placeholders:
            if shape.placeholder_format.type in (
                PP_PLACEHOLDER.TITLE,
                PP_PLACEHOLDER.CENTER_TITLE,
            ):
                title = shape.text
                has_title = bool(title)
                break

        # Find body content
        content = []
        has_content = False
        for shape in slide.placeholders:
            if shape.placeholder_format.type in (
                PP_PLACEHOLDER.BODY,
                PP_PLACEHOLDER.OBJECT,
            ):
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        if p.text:
                            content.append(p.text)
                    if content:
                        has_content = True

        # Inspect all shapes
        has_table = False
        has_chart = False
        has_image = False
        shapes = []

        for shape in slide.shapes:
            shape_info: dict[str, Any] = {
                "name": shape.name,
                "shape_type": str(shape.shape_type),
                "left": shape.left,
                "top": shape.top,
                "width": shape.width,
                "height": shape.height,
            }

            if shape.has_table:
                has_table = True
                tbl = shape.table
                table_data = {
                    "rows": len(tbl.rows),
                    "cols": len(tbl.columns),
                    "headers": [
                        tbl.cell(0, c).text for c in range(len(tbl.columns))
                    ],
                }
                shape_info["table"] = table_data

            if shape.has_chart:
                has_chart = True
                shape_info["chart_type"] = str(shape.chart.chart_type)

            try:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    has_image = True
                    shape_info["image"] = {
                        "width": shape.width,
                        "height": shape.height,
                    }
            except Exception:
                pass

            shapes.append(shape_info)

        # Notes
        notes = ""
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text

        return {
            "slide_number": n,
            "layout": layout_name,
            "title": title,
            "content": content,
            "shapes": shapes,
            "notes": notes,
            "has_title": has_title,
            "has_content": has_content,
            "has_table": has_table,
            "has_chart": has_chart,
            "has_image": has_image,
        }

    def describe_all_slides(self) -> list[dict[str, Any]]:
        """Describe all slides in the presentation.

        Returns:
            List of slide description dicts (same format as describe_slide)

        Example:
            >>> for slide in pres.describe_all_slides():
            ...     print(f"Slide {slide['slide_number']}: {slide['title']}")
        """
        return [
            self.describe_slide(i + 1)
            for i in range(len(self._pptx.slides))
        ]

    # --- Editing methods ---

    def update_slide(
        self,
        n: int,
        *,
        title: str | None = None,
        content: list[str | tuple | dict | list] | None = None,
        levels: list[int] | None = None,
        notes: str | None = None,
    ) -> dict[str, Any]:
        """Update an existing slide. Only provided args are changed.

        Args:
            n: Slide number (1-indexed)
            title: New title (None to keep existing)
            content: New body content (None to keep existing)
            levels: Indent levels for new content
            notes: New speaker notes (None to keep existing)

        Returns:
            Updated slide description dict

        Raises:
            SlideNotFoundError: If slide number is out of range

        Example:
            >>> pres.update_slide(2, title="Updated Title", notes="New notes")
        """
        self._validate_slide_number(n)
        slide = self._pptx.slides[n - 1]

        if title is not None:
            title_ph = self._get_placeholder(slide, PP_PLACEHOLDER.TITLE)
            if title_ph is None:
                title_ph = self._get_placeholder(
                    slide, PP_PLACEHOLDER.CENTER_TITLE
                )
            self._set_text_frame(title_ph, title)

        if content is not None:
            paragraphs = parse_content(content, levels)
            formatted_content, formatted_levels = format_for_py2ppt(paragraphs)
            body_ph = self._get_placeholder(slide, PP_PLACEHOLDER.BODY)
            if body_ph is None:
                body_ph = self._get_placeholder(
                    slide, PP_PLACEHOLDER.OBJECT
                )
            self._set_body_content(
                body_ph, formatted_content, formatted_levels
            )

        if notes is not None:
            self.set_notes(n, notes)

        return self.describe_slide(n)

    def delete_slide(self, n: int) -> int:
        """Delete a slide by number.

        Args:
            n: Slide number (1-indexed)

        Returns:
            New slide count after deletion

        Raises:
            SlideNotFoundError: If slide number is out of range

        Example:
            >>> remaining = pres.delete_slide(3)
        """
        self._validate_slide_number(n)
        sldIdLst = self._pptx.slides._sldIdLst
        rId = sldIdLst[n - 1].rId
        self._pptx.part.drop_rel(rId)
        del sldIdLst[n - 1]
        return len(self._pptx.slides)

    def reorder_slides(self, order: list[int]) -> None:
        """Reorder all slides.

        Args:
            order: List of slide numbers in desired order (1-indexed).
                   Must be a permutation of [1..N].

        Raises:
            InvalidDataError: If order is not a valid permutation

        Example:
            >>> pres.reorder_slides([3, 1, 2])  # Slide 3 becomes first
        """
        count = len(self._pptx.slides)
        if sorted(order) != list(range(1, count + 1)):
            raise InvalidDataError(
                f"Order must be a permutation of slide numbers 1-{count}.",
                suggestion=f"Provide a list containing each number from 1 to {count} exactly once.",
                code="INVALID_SLIDE_ORDER",
            )

        sldIdLst = self._pptx.slides._sldIdLst
        sld_ids = list(sldIdLst)
        for elem in sld_ids:
            sldIdLst.remove(elem)
        for pos in order:
            sldIdLst.append(sld_ids[pos - 1])

    def move_slide(self, from_pos: int, to_pos: int) -> None:
        """Move a single slide to a new position.

        Args:
            from_pos: Current slide number (1-indexed)
            to_pos: Target slide number (1-indexed)

        Raises:
            SlideNotFoundError: If either position is out of range

        Example:
            >>> pres.move_slide(5, 2)  # Move slide 5 to position 2
        """
        self._validate_slide_number(from_pos)
        count = len(self._pptx.slides)
        if not isinstance(to_pos, int) or to_pos < 1 or to_pos > count:
            raise SlideNotFoundError(
                f"Target position {to_pos} is out of range.",
                suggestion=f"Use a position between 1 and {count}.",
                code="SLIDE_NOT_FOUND",
            )
        order = list(range(1, count + 1))
        order.remove(from_pos)
        order.insert(to_pos - 1, from_pos)
        self.reorder_slides(order)

    # --- Content intelligence ---

    def add_content_slides(
        self,
        title: str,
        content: str | list[str | tuple | dict | list],
        *,
        max_bullets: int = 6,
        continuation_suffix: str = " (cont.)",
        levels: list[int] | None = None,
        layout: str | int | None = None,
    ) -> list[int]:
        """Auto-split long content across multiple slides.

        Keeps sub-items (level > 0) grouped with their parent.

        Args:
            title: Slide title (continuation slides get suffix appended)
            content: Content items to split
            max_bullets: Maximum bullets per slide before splitting
            continuation_suffix: Suffix for continuation slide titles
            levels: Optional indent levels for each item
            layout: Layout name, index, or None for auto-selection

        Returns:
            List of slide numbers created

        Example:
            >>> slides = pres.add_content_slides("Key Points", [
            ...     "Point 1", "Detail A", "Detail B",
            ...     "Point 2", "Detail C",
            ...     "Point 3", "Point 4", "Point 5",
            ...     "Point 6", "Point 7", "Point 8",
            ... ], levels=[0, 1, 1, 0, 1, 0, 0, 0, 0, 0, 0], max_bullets=4)
        """
        # Normalize content to list
        if isinstance(content, str):
            content_list: list[str | tuple | dict | list] = [
                line for line in content.split("\n") if line.strip()
            ]
        else:
            content_list = list(content)

        paragraphs = parse_content(content_list, levels)

        # Split into chunks, keeping sub-items with their parent
        chunks: list[list[int]] = []  # each chunk is a list of indices
        current_chunk: list[int] = []

        for idx, para in enumerate(paragraphs):
            if len(current_chunk) >= max_bullets and para.level == 0:
                chunks.append(current_chunk)
                current_chunk = []
            current_chunk.append(idx)

        if current_chunk:
            chunks.append(current_chunk)

        slide_numbers = []
        for i, chunk_indices in enumerate(chunks):
            slide_title = (
                title if i == 0 else f"{title}{continuation_suffix}"
            )
            chunk_content = [content_list[j] for j in chunk_indices]
            chunk_levels = [paragraphs[j].level for j in chunk_indices]
            slide_num = self.add_content_slide(
                slide_title,
                chunk_content,
                levels=chunk_levels,
                layout=layout,
            )
            slide_numbers.append(slide_num)

        return slide_numbers

    def save(self, path: str | Path) -> None:
        """Save the presentation.

        Args:
            path: Output file path

        Example:
            >>> pres.save("output.pptx")
        """
        self._pptx.save(path)

    def __repr__(self) -> str:
        return f"Presentation({self.slide_count} slides, template={self._template.path.name})"
